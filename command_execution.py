# command_execution.py
import os
import subprocess
import sys
import shutil
import threading
import re
import time
import pyautogui
import psutil
import cv2
import io
import soundfile as sf
import moviepy.editor as mp
from moviepy.editor import VideoFileClip, TextClip, CompositeVideoClip
import sqlite3
import logging
import pyperclip
import win32gui
import win32con
import win32api
import win32ui
import win32clipboard
import ctypes
from ctypes import windll
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
import requests
from faster_whisper import WhisperModel
import srt
from bs4 import BeautifulSoup
import socket
import winreg
import platform
import zlib
from cryptography.fernet import Fernet
from cryptography.hazmat.primitives import hashes
from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC
from cryptography.hazmat.primitives.asymmetric import rsa, padding
from cryptography.hazmat.primitives import serialization
import hashlib
import pyotp
from playwright.async_api import async_playwright
from googlesearch import search
import base64
import easyocr
import pyttsx3
import speech_recognition as sr
import PyPDF2
from PIL import Image, ImageFilter, ImageGrab
import yfinance as yf
import speedtest
from functools import wraps, lru_cache
import numpy as np
import datetime
import textwrap
from comtypes import CLSCTX_ALL
from pycaw.pycaw import AudioUtilities, IAudioEndpointVolume
from ctypes import cast, POINTER
import concurrent.futures
from typing import Optional, Union
import string
import secrets
import bcrypt
import librosa
from scipy.signal import lfilter
import pandas as pd
import json
from typing import List, Dict, Any
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches        
from io import BytesIO
import aiohttp
import aiofiles
from contextlib import asynccontextmanager
import tempfile


# Setup logging
logging.basicConfig(filename=r'logs\\main.log', level=logging.INFO,
                    format='%(asctime)s - %(levelname)s - %(message)s')



# Constants
WORKSPACE_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "workspace")
OUTPUT_PATH = os.path.join(WORKSPACE_PATH, "output")
VENV_PATH = os.path.join(WORKSPACE_PATH, "workspacevenv")

def error_handler(func):
    @wraps(func)
    def wrapper(*args, **kwargs):
        try:
            return func(*args, **kwargs)
        except Exception as e:
            error_message = f"Error in {func.__name__}: {str(e)}"
            logging.error(error_message)
            logging.error(f"Exception type: {type(e).__name__}")
            logging.error(f"Exception args: {e.args}")
            return {"status": "error", "error_message": error_message, "exception_type": type(e).__name__, "exception_args": e.args}
    return wrapper

def hashable_lru_cache(maxsize=100):
    def decorator(func):
        cached_func = lru_cache(maxsize=maxsize)(func)

        @wraps(func)
        def wrapper(*args, **kwargs):
            # Convert list arguments to tuples
            new_args = tuple(tuple(arg) if isinstance(arg, list) else arg for arg in args)
            new_kwargs = {k: tuple(v) if isinstance(v, list) else v for k, v in kwargs.items()}
            return cached_func(*new_args, **new_kwargs)

        return wrapper

    return decorator

class CommandExecutionAgent:
    persistent_sessions = {}

    @staticmethod
    def decode_output(raw_output):
        encodings = ['utf-8', 'latin-1']
        for encoding in encodings:
            try:
                return raw_output.decode(encoding)
            except UnicodeDecodeError:
                continue
        return raw_output.decode('utf-8', errors='replace')

    @staticmethod
    @error_handler
    def execute_cmd(command, session_name=None):
        logging.info(f"Executing command: {command} in session: {session_name}")
        try:
            if session_name:
                # Handle persistent session
                if session_name not in CommandExecutionAgent.persistent_sessions:
                    CommandExecutionAgent.persistent_sessions[session_name] = {}
                
                # Execute command in the context of the persistent session
                current_dir = CommandExecutionAgent.persistent_sessions[session_name].get('cwd', os.getcwd())
                full_command = f"cd /d {current_dir} && {command}"
                result = subprocess.run(full_command, shell=True, capture_output=True)
                
                # Update session information
                cwd_result = subprocess.run('cd', shell=True, capture_output=True)
                CommandExecutionAgent.persistent_sessions[session_name]['cwd'] = CommandExecutionAgent.decode_output(cwd_result.stdout).strip()
            else:
                # Non-persistent session
                result = subprocess.run(command, shell=True, capture_output=True)
            
            output = CommandExecutionAgent.decode_output(result.stdout).strip()
            error_output = CommandExecutionAgent.decode_output(result.stderr).strip()
            
            logging.info(f"Command output: {output}")
            if error_output:
                logging.error(f"Command error: {error_output}")
            
            return {
                "status": "success",
                "command": command,
                "output": output,
                "error_output": error_output,
                "session": session_name
            }
        except Exception as e:
            error_message = str(e)
            logging.error(f"Error executing command: {error_message}")
            return {"status": "error", "command": command, "error_message": error_message}

    @staticmethod
    @error_handler
    def close_session(session_name):
        logging.info(f"Closing session: {session_name}")
        try:
            if session_name in CommandExecutionAgent.persistent_sessions:
                del CommandExecutionAgent.persistent_sessions[session_name]
                return {"status": "success", "message": f"Session {session_name} closed"}
            else:
                return {"status": "error", "message": f"Session {session_name} not found"}
        except Exception as e:
            error_message = str(e)
            logging.error(f"Error closing session: {error_message}")
            return {"status": "error", "session_name": session_name, "error_message": error_message}

    @staticmethod
    @error_handler
    def execute_python_code(code, path=None, venv_path=None, run_on_thread=False, visible=False):
        logging.info(f"Executing Python code: {code} in path: {path} using venv: {venv_path}, run_on_thread: {run_on_thread}, visible: {visible}")
        try:
            if venv_path:
                activate_this = os.path.join(venv_path, 'bin', 'activate_this.py')
                exec(open(activate_this).read(), {'__file__': activate_this})
            
            if path:
                os.chdir(path)
            else:
                path = os.getcwd()
                os.chdir(path)
            print(path)
            
            def run_code():
                if visible:
                    # Modify the code to keep the window open until the user manually closes it
                    modified_code = code + "\n\nprint('Press Enter to close the window...')\ninput()"
                    
                    with open('temp_script.py', 'w') as f:
                        f.write(modified_code)
                    
                    if sys.platform.startswith('win'):
                        process = subprocess.Popen(['python', 'temp_script.py'], creationflags=subprocess.CREATE_NEW_CONSOLE)
                    else:
                        process = subprocess.Popen(['python', 'temp_script.py'], start_new_session=True)
                    
                    process.wait()  
                    
                    # Clean up the temporary script file
                    os.remove('temp_script.py')
                else:
                    # Capture the output
                    output = io.StringIO()
                    sys.stdout = output
                    exec(code, {}, {})
                    sys.stdout = sys.__stdout__
                    return output.getvalue().strip()

            if run_on_thread:
                thread = threading.Thread(target=run_code)
                thread.start()
                return {"status": "background", "message": "Code execution started in background thread", "thread_id": thread.ident}
            else:
                result = run_code()
                return {"status": "success", "message": "Code execution completed", "output": result}
        except Exception as e:
            error_message = f"Error executing Python code: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "code": code, "error_message": error_message}

    @staticmethod
    @error_handler
    def pip_install(packages, venv_path=None):
        if isinstance(packages, str):
            packages = [packages]
        results = []
        for package in packages:
            logging.info(f"Installing package: {package}")
            try:
                if venv_path:
                    pip_executable = os.path.join(venv_path, "bin", "pip") if os.name != 'nt' else os.path.join(venv_path, "Scripts", "pip.exe")
                else:
                    pip_executable = "pip"
                subprocess.check_call([pip_executable, "install", package])
                results.append({"status": "success", "package": package, "message": f"Package {package} installed successfully"})
            except subprocess.CalledProcessError as e:
                error_message = f"Failed to install package {package}: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "package": package, "message": error_message})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def run_script(script_path, venv_path=None):
        logging.info(f"Running script: {script_path} using venv: {venv_path}")
        try:
            if venv_path:
                if sys.platform.startswith('win'):
                    python_executable = os.path.join(venv_path, 'Scripts', 'python.exe')
                else:
                    python_executable = os.path.join(venv_path, 'bin', 'python')
            else:
                python_executable = sys.executable

            # Ensure paths are properly escaped
            python_executable = python_executable.replace('\\', '\\\\')
            script_path = script_path.replace('\\', '\\\\')

            # Create a temporary wrapper script
            with tempfile.NamedTemporaryFile(mode='w', suffix='.py', delete=False) as temp_file:
                temp_file.write(f"""
import subprocess
import sys

# Run the original script
result = subprocess.run([r'{python_executable}', r'{script_path}'], capture_output=True, text=True)

# Print the output
print(result.stdout)
print(result.stderr, file=sys.stderr)

# Write output to files for later retrieval
with open('script_output.txt', 'w') as out_file:
    out_file.write(result.stdout)
with open('script_error.txt', 'w') as err_file:
    err_file.write(result.stderr)

input('Press Enter to close the window...')
""")
                temp_wrapper_path = temp_file.name

            # Run the wrapper script in a new console
            if sys.platform.startswith('win'):
                process = subprocess.Popen([python_executable, temp_wrapper_path], creationflags=subprocess.CREATE_NEW_CONSOLE)
            else:
                process = subprocess.Popen([python_executable, temp_wrapper_path], start_new_session=True)

            process.wait()

            # Retrieve the output from the files
            output = ""
            error_output = ""
            try:
                with open('script_output.txt', 'r') as out_file:
                    output = out_file.read()
                with open('script_error.txt', 'r') as err_file:
                    error_output = err_file.read()
            except FileNotFoundError:
                logging.warning("Output files not found. The script might have failed to execute.")

            # Clean up temporary files
            os.remove(temp_wrapper_path)
            try:
                os.remove('script_output.txt')
                os.remove('script_error.txt')
            except FileNotFoundError:
                pass

            if error_output:
                logging.error(f"Script error: {error_output}")
                return {"status": "error", "script_path": script_path, "error_output": error_output}
            else:
                logging.info(f"Script output: {output}")
                return {"status": "success", "script_path": script_path, "output": output}
        except Exception as e:
            error_message = f"Error running script: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "script_path": script_path, "error_output": error_message}

class FileOperationsAgent:
    @staticmethod
    @error_handler
    def create_file(paths, content=''):
        if isinstance(paths, str):
            paths = [paths]
        results = []
        for path in paths:
            logging.info(f"Creating file at: {path} with content: {content}")
            try:
                with open(path, 'w') as f:
                    f.write(content)
                success_message = f"File created at {path}"
                logging.info(success_message)
                results.append({"status": "success", "path": path, "message": success_message})
            except Exception as e:
                error_message = f"Error creating file: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "path": path, "error_message": error_message})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def read_file(paths):
        if isinstance(paths, str):
            paths = [paths]
        results = []
        for path in paths:
            logging.info(f"Reading file at: {path}")
            try:
                with open(path, 'r') as f:
                    content = f.read()
                logging.info(f"Content read from file: {content}")
                results.append({"status": "success", "path": path, "content": content})
            except Exception as e:
                error_message = f"Error reading file: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "path": path, "error_message": error_message})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def delete_file(paths):
        if isinstance(paths, str):
            paths = [paths]
        results = []
        for path in paths:
            logging.info(f"Deleting file at: {path}")
            try:
                os.remove(path)
                success_message = f"File deleted: {path}"
                logging.info(success_message)
                results.append({"status": "success", "path": path, "message": success_message})
            except Exception as e:
                error_message = f"Error deleting file: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "path": path, "error_message": error_message})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def copy_file(src_dst_pairs):
        if isinstance(src_dst_pairs[0], str):
            src_dst_pairs = [src_dst_pairs]
        results = []
        for src, dst in src_dst_pairs:
            logging.info(f"Copying file from {src} to {dst}")
            try:
                shutil.copy2(src, dst)
                success_message = f"File copied from {src} to {dst}"
                logging.info(success_message)
                results.append({"status": "success", "src": src, "dst": dst, "message": success_message})
            except Exception as e:
                error_message = f"Error copying file: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "src": src, "dst": dst, "error_message": error_message})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def move_file(src_dst_pairs):
        if isinstance(src_dst_pairs[0], str):
            src_dst_pairs = [src_dst_pairs]
        results = []
        for src, dst in src_dst_pairs:
            logging.info(f"Moving file from {src} to {dst}")
            try:
                shutil.move(src, dst)
                success_message = f"File moved from {src} to {dst}"
                logging.info(success_message)
                results.append({"status": "success", "src": src, "dst": dst, "message": success_message})
            except Exception as e:
                error_message = f"Error moving file: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "src": src, "dst": dst, "error_message": error_message})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def rename_file(src_dst_pairs):
        if isinstance(src_dst_pairs[0], str):
            src_dst_pairs = [src_dst_pairs]
        results = []
        for src, dst in src_dst_pairs:
            logging.info(f"Renaming file from {src} to {dst}")
            try:
                os.rename(src, dst)
                success_message = f"File renamed from {src} to {dst}"
                logging.info(success_message)
                results.append({"status": "success", "src": src, "dst": dst, "message": success_message})
            except Exception as e:
                error_message = f"Error renaming file: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "src": src, "dst": dst, "error_message": error_message})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def append_to_file(path_content_pairs):
        if isinstance(path_content_pairs[0], str):
            path_content_pairs = [path_content_pairs]
        results = []
        for path, content in path_content_pairs:
            logging.info(f"Appending to file at: {path} with content: {content}")
            try:
                with open(path, 'a') as f:
                    f.write(content)
                success_message = f"Content appended to file at {path}"
                logging.info(success_message)
                results.append({"status": "success", "path": path, "message": success_message})
            except Exception as e:
                error_message = f"Error appending to file: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "path": path, "error_message": error_message})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def file_exists(paths):
        if isinstance(paths, str):
            paths = [paths]
        results = []
        for path in paths:
            logging.info(f"Checking existence of file at: {path}")
            try:
                exists = os.path.exists(path)
                message = f"File exists: {exists}" if exists else f"File does not exist: {path}"
                logging.info(message)
                results.append({"status": "success", "path": path, "exists": exists, "message": message})
            except Exception as e:
                error_message = f"Error checking file existence: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "path": path, "error_message": error_message})
        return results if len(results) > 1 else results[0]


    @staticmethod
    @error_handler
    def list_directory(path):
        logging.info(f"Listing directory at: {path}")
        try:
            items = os.listdir(path)
            logging.info(f"Items in directory: {items}")
            return {"status": "success", "path": path, "items": items}
        except Exception as e:
            error_message = f"Error listing directory: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "path": path, "error_message": error_message}


    @staticmethod
    @error_handler
    def create_directory(paths):
        if isinstance(paths, str):
            paths = [paths]
        results = []
        for path in paths:
            logging.info(f"Creating directory at: {path}")
            try:
                os.makedirs(path, exist_ok=True)
                success_message = f"Directory created: {path}"
                logging.info(success_message)
                results.append({"status": "success", "path": path, "message": success_message})
            except Exception as e:
                error_message = f"Error creating directory: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "path": path, "error_message": error_message})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def delete_directory(paths):
        if isinstance(paths, str):
            paths = [paths]
        results = []
        for path in paths:
            logging.info(f"Deleting directory at: {path}")
            try:
                shutil.rmtree(path)
                success_message = f"Directory deleted: {path}"
                logging.info(success_message)
                results.append({"status": "success", "path": path, "message": success_message})
            except Exception as e:
                error_message = f"Error deleting directory: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "path": path, "error_message": error_message})
        return results if len(results) > 1 else results[0]


    @staticmethod
    @error_handler
    def search_files(directory, pattern):
        logging.info(f"Searching for files in {directory} matching pattern: {pattern}")
        try:
            matches = []
            for root, _, filenames in os.walk(directory):
                for filename in filenames:
                    if re.search(pattern, filename):
                        matches.append(os.path.join(root, filename))
            logging.info(f"Found matches: {matches}")
            return {"status": "success", "directory": directory, "pattern": pattern, "matches": matches}
        except Exception as e:
            error_message = f"Error searching files: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "directory": directory, "pattern": pattern, "error_message": error_message}

    @staticmethod
    @error_handler
    def get_file_info(paths):
        if isinstance(paths, str):
            paths = [paths]
        results = []
        for path in paths:
            logging.info(f"Getting file info for: {path}")
            try:
                info = {
                    'size': os.path.getsize(path),
                    'created': time.ctime(os.path.getctime(path)),
                    'modified': time.ctime(os.path.getmtime(path)),
                    'accessed': time.ctime(os.path.getatime(path))
                }
                logging.info(f"File info: {info}")
                results.append({"status": "success", "path": path, "info": info})
            except Exception as e:
                error_message = f"Error getting file info: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "path": path, "error_message": error_message})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def monitor_directory(path, duration):
        logging.info(f"Monitoring directory at: {path} for {duration} seconds")
        try:
            class Handler(FileSystemEventHandler):
                def __init__(self):
                    self.events = []

                def on_any_event(self, event):
                    self.events.append(f"{event.event_type}: {event.src_path}")

            event_handler = Handler()
            observer = Observer()
            observer.schedule(event_handler, path, recursive=True)
            observer.start()
            time.sleep(duration)
            observer.stop()
            observer.join()
            logging.info(f"Directory events: {event_handler.events}")
            return {"status": "success", "path": path, "duration": duration, "events": event_handler.events}
        except Exception as e:
            error_message = f"Error monitoring directory: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "path": path, "duration": duration, "error_message": error_message}


    @staticmethod
    @error_handler
    def get_file_hash(file_paths, algorithm='sha256'):
        if isinstance(file_paths, str):
            file_paths = [file_paths]
        results = []
        for file_path in file_paths:
            logging.info(f"Computing {algorithm} hash for file: {file_path}")
            try:
                hash_object = hashlib.new(algorithm)
                with open(file_path, 'rb') as file:
                    for chunk in iter(lambda: file.read(4096), b''):
                        hash_object.update(chunk)
                file_hash = hash_object.hexdigest()
                logging.info(f"Computed hash: {file_hash}")
                results.append({"status": "success", "file_path": file_path, "algorithm": algorithm, "hash": file_hash})
            except Exception as e:
                error_message = f"Error computing hash: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "file_path": file_path, "algorithm": algorithm, "error_message": error_message})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def compress_file(input_output_pairs):
        if isinstance(input_output_pairs[0], str):
            input_output_pairs = [input_output_pairs]
        results = []
        for input_file, output_file in input_output_pairs:
            logging.info(f"Compressing file from {input_file} to {output_file}")
            try:
                with open(input_file, 'rb') as f_in, open(output_file, 'wb') as f_out:
                    f_out.write(zlib.compress(f_in.read()))
                success_message = f"File compressed to {output_file}"
                logging.info(success_message)
                results.append({"status": "success", "input_file": input_file, "output_file": output_file, "message": success_message})
            except Exception as e:
                error_message = f"Error compressing file: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "input_file": input_file, "output_file": output_file, "error_message": error_message})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def decompress_file(input_output_pairs):
        if isinstance(input_output_pairs[0], str):
            input_output_pairs = [input_output_pairs]
        results = []
        for input_file, output_file in input_output_pairs:
            logging.info(f"Decompressing file from {input_file} to {output_file}")
            try:
                with open(input_file, 'rb') as f_in, open(output_file, 'wb') as f_out:
                    f_out.write(zlib.decompress(f_in.read()))
                success_message = f"File decompressed to {output_file}"
                logging.info(success_message)
                results.append({"status": "success", "input_file": input_file, "output_file": output_file, "message": success_message})
            except Exception as e:
                error_message = f"Error decompressing file: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "input_file": input_file, "output_file": output_file, "error_message": error_message})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def merge_pdfs(pdf_files, output_file):
        logging.info(f"Merging PDFs: {pdf_files} into {output_file}")
        try:
            merger = PyPDF2.PdfMerger()
            for pdf in pdf_files:
                merger.append(pdf)
            merger.write(output_file)
            merger.close()
            success_message = f"PDFs merged into {output_file}"
            logging.info(success_message)
            return {"status": "success", "pdf_files": pdf_files, "output_file": output_file, "message": success_message}
        except Exception as e:
            error_message = f"Error merging PDFs: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "pdf_files": pdf_files, "output_file": output_file, "error_message": error_message}


    @staticmethod
    @error_handler
    def split_pdf(input_file, start_page, end_page, output_file):
        logging.info(f"Splitting PDF: {input_file} from page {start_page} to {end_page} into {output_file}")
        try:
            with open(input_file, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                writer = PyPDF2.PdfWriter()
                for page in range(start_page - 1, min(end_page, len(reader.pages))):
                    writer.add_page(reader.pages[page])
                with open(output_file, 'wb') as output:
                    writer.write(output)
            success_message = f"PDF split and saved to {output_file}"
            logging.info(success_message)
            return {"status": "success", "input_file": input_file, "start_page": start_page, "end_page": end_page, "output_file": output_file, "message": success_message}
        except Exception as e:
            error_message = f"Error splitting PDF: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "input_file": input_file, "start_page": start_page, "end_page": end_page, "output_file": output_file, "error_message": error_message}


class WebInteractionAgent:
    def __init__(self, headless=True):
        self.headless = headless
        self.browser = None
        self.context = None
        self.page = None

    async def initialize(self):
        self.playwright = await async_playwright().start()
        self.browser = await self.playwright.chromium.launch(headless=self.headless)
        self.context = await self.browser.new_context()
        self.page = await self.context.new_page()
    
    @asynccontextmanager
    async def managed_browser(self):
        await self.initialize()
        try:
            yield self
        finally:
            await self.close()

    async def close(self):
        await self.page.close()
        await self.context.close()
        await self.browser.close()
        await self.playwright.stop()
        
    @error_handler
    async def get_all_elements(self, url):
        async with self.managed_browser():
            logging.info(f"Retrieving all elements from {url}")
            await self.page.goto(url)
            elements = await self.page.query_selector_all('*')
            element_data = []
            for element in elements:
                try:
                    attrs = await element.evaluate('''
                        (el) => {
                            const attributes = {};
                            for (let attr of el.attributes) {
                                attributes[attr.name] = attr.value;
                            }
                            return attributes;
                        }
                    ''')
                    element_data.append({
                        'tag_name': await element.evaluate('(el) => el.tagName.toLowerCase()'),
                        'text': await element.inner_text() if await element.evaluate('(el) => el instanceof HTMLElement') else '',
                        'attributes': attrs
                    })
                except Exception as e:
                    logging.warning(f"Error processing element: {str(e)}")
            logging.info(f"Retrieved {len(elements)} elements from {url}")
            return {"status": "success", "url": url, "elements": element_data}

    @error_handler
    async def interact_with_element(self, url, element_selector, action, text=None):
        async with self.managed_browser():
            logging.info(f"Interacting with element at {url} using {action} action")
            await self.page.goto(url)
            element = await self.page.query_selector(element_selector)
            
            if element is None:
                return {"status": "error", "message": f"Element with selector '{element_selector}' not found"}
            
            if action == 'click':
                await element.click()
            elif action == 'enter_text':
                if text is None:
                    raise ValueError("Text must be provided for entering text")
                await element.fill(text)
            elif action == 'hover':
                await element.hover()
            elif action == 'submit':
                await element.evaluate('form => form.submit()')
            else:
                raise ValueError("Unsupported action")
            
            success_message = f"Action {action} performed on element {element_selector} at {url}"
            logging.info(success_message)
            return {"status": "success", "url": url, "element_selector": element_selector, "action": action, "message": success_message}

    @error_handler
    async def web_search(self, query, num_results=5):
        async with self.managed_browser():
            logging.info(f"Performing web search for query: {query}")
            try:
                results = search(query, num_results=num_results)
                result_data = [result for result in results]
                success_message = f"Search results retrieved for query: {query}"
                logging.info(success_message)
                return {"status": "success", "query": query, "results": result_data}
            except Exception as e:
                error_message = f"Error performing web search: {str(e)}"
                logging.error(error_message)
                return {"status": "error", "error_message": error_message}
        
    @error_handler
    async def check_url_status(self, urls):
        async with self.managed_browser():
            if isinstance(urls, str):
                urls = [urls]
            results = []
            for url in urls:
                logging.info(f"Checking URL status for {url}")
                try:
                    response = requests.head(url)
                    response.raise_for_status()
                    status_code = response.status_code
                    logging.info(f"URL {url} status code: {status_code}")
                    results.append({"status": "success", "url": url, "status_code": status_code})
                except Exception as e:
                    error_message = f"Error checking URL status: {str(e)}"
                    logging.error(error_message)
                    results.append({"status": "error", "url": url, "error_message": error_message})
            return results if len(results) > 1 else results[0]
    
    @error_handler
    async def extract_attributes(self, url, css_selector, attribute):
        async with self.managed_browser():
            logging.info(f"Extracting attribute {attribute} from elements at {url} with selector {css_selector}")
            await self.page.goto(url)
            elements = await self.page.query_selector_all(css_selector)
            attributes = [await element.get_attribute(attribute) for element in elements if await element.get_attribute(attribute)]
            logging.info(f"Extracted attributes: {attributes}")
            return {"status": "success", "url": url, "css_selector": css_selector, "attribute": attribute, "attributes": attributes}

    @error_handler
    async def download_all_files(self, url, file_extension, destination_folder):
        async with self.managed_browser():
            logging.info(f"Downloading all files with extension {file_extension} from {url} to {destination_folder}")
            await self.page.goto(url)
            soup = BeautifulSoup(await self.page.content(), 'html.parser')
            file_links = [await link.get('href') for link in await soup.find_all('a', href=True) if link['href'].endswith(file_extension)]
            
            downloaded_files = []
            os.makedirs(destination_folder, exist_ok=True)
            for link in file_links:
                file_url = link if link.startswith('http') else url + '/' + link
                file_name = os.path.basename(link)
                file_path = os.path.join(destination_folder, file_name)
                try:
                    download_response = requests.get(file_url, stream=True)
                    download_response.raise_for_status()
                    with open(file_path, 'wb') as file:
                        for chunk in download_response.iter_content(chunk_size=8192):
                            file.write(chunk)
                    downloaded_files.append(file_path)
                    logging.info(f"Downloaded file: {file_path}")
                except Exception as e:
                    logging.error(f"Failed to download {file_url}: {str(e)}")
            
            success_message = f"Downloaded files: {downloaded_files}"
            logging.info(success_message)
            return {"status": "success", "url": url, "file_extension": file_extension, "downloaded_files": downloaded_files, "message": success_message}

    @error_handler
    async def web_request(self, urls, method='get', params=None, headers=None, data=None, json=None):
        async with self.managed_browser():
            if isinstance(urls, str):
                urls = [urls]
            results = []
            for url in urls:
                logging.info(f"Making {method.upper()} request to {url}")
                try:
                    if method.lower() == 'get':
                        response = requests.get(url, params=params, headers=headers)
                    elif method.lower() == 'post':
                        response = requests.post(url, params=params, headers=headers, data=data, json=json)
                    else:
                        error_message = f"Unsupported method: {method}"
                        logging.error(error_message)
                        results.append({"status": "error", "method": method, "url": url, "error_message": error_message})
                        continue

                    response.raise_for_status()
                    response_text = response.text
                    logging.info(f"Response received from {url}")
                    results.append({"status": "success", "method": method, "url": url, "response": response_text})
                except Exception as e:
                    error_message = f"Error making web request: {str(e)}"
                    logging.error(error_message)
                    results.append({"status": "error", "method": method, "url": url, "error_message": error_message})
            return results if len(results) > 1 else results[0]

    @error_handler
    async def download_file(self, url_destination_pairs):
        async with self.managed_browser():
            if isinstance(url_destination_pairs[0], str):
                url_destination_pairs = [url_destination_pairs]
            results = []
            async with aiohttp.ClientSession() as session:
                for url, destination in url_destination_pairs:
                    logging.info(f"Downloading file from {url} to {destination}")
                    try:
                        async with session.get(url) as response:
                            response.raise_for_status()
                            async with aiofiles.open(destination, 'wb') as file:
                                async for chunk in response.content.iter_chunked(8192):
                                    await file.write(chunk)
                        success_message = f"File downloaded from {url} to {destination}"
                        logging.info(success_message)
                        results.append({"status": "success", "url": url, "destination": destination, "message": success_message})
                    except Exception as e:
                        error_message = f"Error downloading file: {str(e)}"
                        logging.error(error_message)
                        results.append({"status": "error", "url": url, "destination": destination, "error_message": error_message})
            return results if len(results) > 1 else results[0]

    @error_handler
    async def scrape_website(self, url_selector_pairs):
        async with self.managed_browser():
            if isinstance(url_selector_pairs[0], str):
                url_selector_pairs = [url_selector_pairs]
            results = []
            async with aiohttp.ClientSession() as session:
                for url, css_selector in url_selector_pairs:
                    logging.info(f"Scraping website {url} with selector {css_selector}")
                    try:
                        async with session.get(url) as response:
                            response.raise_for_status()
                            html = await response.text()
                            soup = BeautifulSoup(html, 'html.parser')
                            elements = soup.select(css_selector)
                            content = [element.text for element in elements]
                            logging.info(f"Scraped content: {content}")
                            results.append({"status": "success", "url": url, "css_selector": css_selector, "content": content})
                    except Exception as e:
                        error_message = f"Error scraping website: {str(e)}"
                        logging.error(error_message)
                        results.append({"status": "error", "url": url, "css_selector": css_selector, "error_message": error_message})
            return results if len(results) > 1 else results[0]


class SystemManagementAgent:
    @staticmethod
    @error_handler
    def system_health_check():
        logging.info("Performing system health check")
        try:
            health = {
                'cpu_percent': psutil.cpu_percent(interval=1),
                'memory_percent': psutil.virtual_memory().percent,
                'disk_percent': psutil.disk_usage('/').percent,
                'battery': psutil.sensors_battery().percent if psutil.sensors_battery() else None,
            }
            if hasattr(psutil, 'sensors_temperatures'):
                health['cpu_temperature'] = psutil.sensors_temperatures().get('coretemp', [None])[0].current if psutil.sensors_temperatures() else None
            logging.info(f"System health: {health}")
            return {"status": "success", "health": health}
        except Exception as e:
            error_message = f"Error performing system health check: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}

    @staticmethod
    @error_handler
    def clean_system():
        logging.info("Initiating system cleanup")
        try:
            if platform.system() == 'Windows':
                cleanmgr_path = os.getenv('CLEANMGR_PATH', r'C:\Windows\System32\cleanmgr.exe')
                if os.path.exists(cleanmgr_path):
                    command = f'"{cleanmgr_path}" /sagerun:1'
                    result = CommandExecutionAgent.execute_cmd(command)
                else:
                    raise FileNotFoundError(f"cleanmgr.exe not found at {cleanmgr_path}")
            elif platform.system() == 'Linux':
                result = CommandExecutionAgent.execute_cmd('sudo apt-get autoremove -y && sudo apt-get autoclean -y')
            elif platform.system() == 'Darwin':
                result = CommandExecutionAgent.execute_cmd('brew cleanup')
            else:
                return {"status": "error", "message": "System cleanup not supported on this platform"}
            
            if result['status'] == 'success':
                success_message = "System cleanup initiated"
                logging.info(success_message)
                return {"status": "success", "message": success_message, "output": result['output']}
            else:
                return {"status": "error", "error_message": result['error_output']}
        except Exception as e:
            error_message = f"Error performing system cleanup: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}
            
    @staticmethod
    @error_handler
    def shutdown_system(delay=3):
        logging.info(f"Shutting down system in {delay} seconds")
        try:
            if platform.system() == 'Windows':
                # Use the full path to the shutdown command
                shutdown_path = os.getenv('SHUTDOWN_PATH', r'C:\Windows\System32\shutdown.exe')
                if os.path.exists(shutdown_path):
                    command = f'"{shutdown_path}" /s /t {delay}'
                else:
                    raise FileNotFoundError(f"Shutdown executable not found at {shutdown_path}")
            elif platform.system() in ['Linux', 'Darwin']:
                command = f'sudo shutdown -h +{delay // 60}'
            else:
                raise ValueError("Unsupported platform")
            
            logging.info(f"Executing command: {command}")
            result = CommandExecutionAgent.execute_cmd(command)
            
            if result['status'] == 'success':
                success_message = f"System will shut down in {delay} seconds"
                logging.info(success_message)
                return {"status": "success", "message": success_message, "output": result['output']}
            else:
                raise Exception(f"Command execution failed: {result['error_output']}")
        except Exception as e:
            error_message = f"Error shutting down system: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}

    @staticmethod
    @error_handler
    def restart_system(delay=3):
        logging.info(f"Restarting system in {delay} seconds")
        try:
            if platform.system() == 'Windows':
                # Use the full path to the shutdown command
                shutdown_path = os.getenv('SHUTDOWN_PATH', r'C:\Windows\System32\shutdown.exe')
                if os.path.exists(shutdown_path):
                    command = f'"{shutdown_path}" /r /t {delay}'
                else:
                    raise FileNotFoundError(f"Shutdown executable not found at {shutdown_path}")
            elif platform.system() in ['Linux', 'Darwin']:
                command = f'sudo shutdown -r +{delay // 60}'
            else:
                raise ValueError("Unsupported platform")
            
            logging.info(f"Executing command: {command}")
            result = CommandExecutionAgent.execute_cmd(command)
            
            if result['status'] == 'success':
                success_message = f"System will restart in {delay} seconds"
                logging.info(success_message)
                return {"status": "success", "message": success_message, "output": result['output']}
            else:
                raise Exception(f"Command execution failed: {result['error_output']}")
        except Exception as e:
            error_message = f"Error restarting system: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}

    @staticmethod
    @error_handler
    def sleep_system():
        logging.info("Putting system to sleep")
        try:
            if platform.system() == 'Windows':
                rundll32_path = os.getenv('RUNDLL32_PATH', r'C:\Windows\System32\rundll32.exe')
                if os.path.exists(rundll32_path):
                    command = f'"{rundll32_path}" powrprof.dll,SetSuspendState 0,1,0'
                else:
                    raise FileNotFoundError(f"rundll32.exe not found at {rundll32_path}")
            elif platform.system() == 'Linux':
                command = 'systemctl suspend'
            elif platform.system() == 'Darwin':
                command = 'pmset sleepnow'
            else:
                raise ValueError("Unsupported platform")
            
            logging.info(f"Executing command: {command}")
            result = CommandExecutionAgent.execute_cmd(command)
            
            if result['status'] == 'success':
                success_message = "System is entering sleep mode"
                logging.info(success_message)
                return {"status": "success", "message": success_message, "output": result['output']}
            else:
                raise Exception(f"Command execution failed: {result['error_output']}")
        except Exception as e:
            error_message = f"Error putting system to sleep: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}

    @staticmethod
    @error_handler
    def get_system_info():
        logging.info("Retrieving system information")
        try:
            info = {
                'os': platform.system(),
                'os_version': platform.version(),
                'cpu': platform.processor(),
                'cpu_cores': psutil.cpu_count(logical=False),
                'cpu_threads': psutil.cpu_count(logical=True),
                'ram': psutil.virtual_memory().total,
                'disk': psutil.disk_usage('/').total,
                'ip_address': None
            }
            
            try:
                ip_address = None
                interfaces = psutil.net_if_addrs()
                for interface, addrs in interfaces.items():
                    for addr in addrs:
                        if addr.family == socket.AF_INET and not addr.address.startswith("127."):
                            ip_address = addr.address
                            break
                    if ip_address:
                        break
                
                if not ip_address:
                    ip_address = next((addr.address for interface, addrs in interfaces.items() 
                                    for addr in addrs if addr.family == socket.AF_INET), None)
                
                info['ip_address'] = ip_address
            except Exception as e:
                logging.warning(f"Error retrieving IP address: {str(e)}")
            
            logging.info(f"System information: {info}")
            return {"status": "success", "info": info}
        except Exception as e:
            error_message = f"Error retrieving system information: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}

    @staticmethod
    @error_handler
    def get_battery_status():
        logging.info("Retrieving battery status")
        try:
            battery = psutil.sensors_battery()
            if battery:
                status = {
                    'percent': battery.percent,
                    'power_plugged': battery.power_plugged,
                    'time_left': battery.secsleft
                }
            else:
                status = None
            logging.info(f"Battery status: {status}")
            return {"status": "success", "battery_status": status}
        except Exception as e:
            error_message = f"Error retrieving battery status: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}

    @staticmethod
    @error_handler
    def list_installed_software():
        logging.info("Listing installed software")
        software_list = []
        try:
            if platform.system() == 'Windows':
                key = winreg.OpenKey(winreg.HKEY_LOCAL_MACHINE, r"SOFTWARE\Microsoft\Windows\CurrentVersion\Uninstall")
                for i in range(0, winreg.QueryInfoKey(key)[0]):
                    try:
                        subkey_name = winreg.EnumKey(key, i)
                        subkey = winreg.OpenKey(key, subkey_name)
                        display_name = winreg.QueryValueEx(subkey, "DisplayName")[0]
                        software_list.append(display_name)
                    except WindowsError:
                        pass
            elif platform.system() == 'Linux':
                installed_packages = subprocess.check_output(['dpkg', '--get-selections']).decode('utf-8').split('\n')
                software_list = [package.split('\t')[0] for package in installed_packages if package]
            elif platform.system() == 'Darwin':
                installed_packages = subprocess.check_output(['brew', 'list']).decode('utf-8').split('\n')
                software_list = [package for package in installed_packages if package]
            else:
                raise ValueError("Unsupported platform")
            logging.info(f"Installed software: {software_list}")
            return {"status": "success", "installed_software": software_list}
        except Exception as e:
            error_message = f"Error listing installed software: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}

    @staticmethod
    @error_handler
    def get_startup_programs():
        logging.info("Retrieving startup programs")
        startup_list = []
        try:
            if platform.system() == 'Windows':
                startup_folder = os.path.join(os.getenv('APPDATA'), 'Microsoft', 'Windows', 'Start Menu', 'Programs', 'Startup')
                startup_list = os.listdir(startup_folder)
            elif platform.system() == 'Linux':
                startup_list = os.listdir('/etc/init.d/')
            elif platform.system() == 'Darwin':
                startup_list = os.listdir('/Library/StartupItems/')
            else:
                raise ValueError("Unsupported platform")
            logging.info(f"Startup programs: {startup_list}")
            return {"status": "success", "startup_programs": startup_list}
        except Exception as e:
            error_message = (f"Error retrieving startup programs: {str(e)}")
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}

    @staticmethod
    @error_handler
    def optimize_memory():
        logging.info("Optimizing memory usage")
        try:
            if platform.system() == 'Windows':
                powershell_exe = os.getenv('POWERSHELL_EXE_PATH', r'C:\Windows\System32\WindowsPowerShell\v1.0\powershell.exe')
                ipconfig_exe = os.getenv('IPCONFIG_EXE_PATH', r'C:\Windows\System32\ipconfig.exe')

                commands = [
                    f'"{powershell_exe}" -Command "Clear-RecycleBin -Force -ErrorAction SilentlyContinue"',
                    f'"{ipconfig_exe}" /flushdns',
                    f'"{powershell_exe}" -Command "Clear-DNSClientCache"',
                    f'"{powershell_exe}" -Command "Stop-Service -Name SysMain -Force; Start-Service -Name SysMain"',
                    f'"{powershell_exe}" -Command "Get-Process | Sort-Object -Property WS -Descending | Select-Object -First 5 | ForEach-Object {{ Write-Output \\"Top memory consumer: $($_.Name) (Current memory: $([Math]::Round($_.WS / 1MB, 2)) MB)\\" }}"',
                    f'"{powershell_exe}" -Command "[System.GC]::Collect()"',
                ]
                outputs = []
                for command in commands:
                    result = CommandExecutionAgent.execute_cmd(command)
                    if result['status'] == 'success':
                        if result['output']:
                            outputs.append(result['output'])
                        if result['error_output']:
                            outputs.append(f"Warning: {result['error_output']}")
                    else:
                        outputs.append(f"Error executing command: {result.get('error_message', 'Unknown error')}")
                output = "\n".join(outputs)
            elif platform.system() == 'Linux':
                command = 'sync && echo 1 | sudo tee /proc/sys/vm/drop_caches'
                result = CommandExecutionAgent.execute_cmd(command)
                output = result['output']
            elif platform.system() == 'Darwin':
                command = 'sudo purge'
                result = CommandExecutionAgent.execute_cmd(command)
                output = result['output']
            else:
                raise ValueError("Unsupported platform")
            
            success_message = "Memory optimization complete"
            logging.info(success_message)
            return {"status": "success", "message": success_message, "output": output}
        except Exception as e:
            error_message = f"Error optimizing memory: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}

    @staticmethod
    @error_handler
    def update_system():
        logging.info("Updating system")
        try:
            if platform.system() == 'Windows':
                powershell_path = os.getenv('POWERSHELL_EXE_PATH', r'C:\Windows\System32\WindowsPowerShell\v1.0\powershell.exe')
                if os.path.exists(powershell_path):
                    command = f'"{powershell_path}" -Command "Install-Module -Name PSWindowsUpdate -Force; Get-WindowsUpdate -Install -AcceptAll"'
                else:
                    raise FileNotFoundError(f"PowerShell not found at {powershell_path}")
            elif platform.system() == 'Linux':
                command = 'sudo apt-get update && sudo apt-get upgrade -y'
            elif platform.system() == 'Darwin':
                command = 'softwareupdate -i -a'
            else:
                raise ValueError("Unsupported platform")
            
            result = CommandExecutionAgent.execute_cmd(command)
            
            if result['status'] == 'success':
                success_message = "System update complete"
                logging.info(success_message)
                return {"status": "success", "message": success_message, "output": result['output']}
            else:
                return {"status": "error", "error_message": result['error_output']}
        except Exception as e:
            error_message = f"Error updating system: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}

    @staticmethod
    @error_handler
    def defragment_disk(drive_letters='C:'):
        if isinstance(drive_letters, str):
            drive_letters = [drive_letters]
        results = []
        for drive_letter in drive_letters:
            logging.info(f"Defragmenting disk {drive_letter}")
            try:
                if platform.system() == 'Windows':
                    defrag_path = os.getenv('DEFRAG_PATH', r'C:\Windows\System32\defrag.exe')
                    if os.path.exists(defrag_path):
                        command = f'"{defrag_path}" {drive_letter} /O'
                    else:
                        raise FileNotFoundError(f"Defrag executable not found at {defrag_path}")
                    
                    logging.info(f"Executing command: {command}")
                    result = CommandExecutionAgent.execute_cmd(command)
                    
                    if result['status'] == 'success':
                        success_message = f"Defragmentation of {drive_letter} complete"
                        logging.info(success_message)
                        results.append({"status": "success", "drive_letter": drive_letter, "message": success_message, "output": result['output']})
                    else:
                        raise Exception(f"Command execution failed: {result['error_output']}")
                else:
                    raise ValueError("Defragmentation is not supported on this platform")
            except Exception as e:
                error_message = f"Error defragmenting disk {drive_letter}: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "drive_letter": drive_letter, "error_message": error_message})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def check_disk_health(disks='C:'):
        if isinstance(disks, str):
            disks = [disks]
        results = []
        for disk in disks:
            logging.info(f"Checking disk health for {disk}")
            try:
                if platform.system() == 'Windows':
                    result = subprocess.check_output(['chkdsk', disk], stderr=subprocess.STDOUT, universal_newlines=True, encoding='latin-1')
                    health = "Healthy" if "Windows has scanned the file system and found no problems" in result else "Potential issues found"
                elif platform.system() == 'Linux':
                    result = subprocess.check_output(['sudo', 'fsck', '-n', disk], shell=True, universal_newlines=True)
                    health = "Healthy" if "clean" in result else "Potential issues found"
                elif platform.system() == 'Darwin':
                    result = subprocess.check_output(['diskutil', 'verifyVolume', disk], shell=True, universal_newlines=True)
                    health = "Healthy" if "appears to be OK" in result else "Potential issues found"
                else:
                    raise ValueError("Unsupported platform")
                logging.info(f"Disk health for {disk}: {health}")
                results.append({"status": "success", "disk": disk, "health": health, "output": result})
            except subprocess.CalledProcessError as e:
                error_message = f"Error checking disk health for {disk}: {e.output}"
                logging.error(error_message)
                results.append({"status": "error", "disk": disk, "error_message": error_message})
            except Exception as e:
                error_message = f"Error checking disk health for {disk}: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "disk": disk, "error_message": error_message})
        return results if len(results) > 1 else results[0]
        
    @staticmethod
    @error_handler
    def clear_temp_files():
        logging.info("Clearing temporary files")
        try:
            if platform.system() == 'Windows':
                cmd_path = os.getenv('CMD_PATH', r'C:\Windows\System32\cmd.exe')
                command = f'{cmd_path} /c del /q/f/s %TEMP%\*'
            elif platform.system() == 'Linux':
                command = 'sudo rm -rf /tmp/*'
            elif platform.system() == 'Darwin':
                command = 'sudo rm -rf /private/var/tmp/*'
            else:
                raise ValueError("Unsupported platform")
            
            result = CommandExecutionAgent.execute_cmd(command)
            
            if result['status'] == 'success':
                success_message = "Temporary files cleared"
                logging.info(success_message)
                return {"status": "success", "message": success_message, "output": "\n".join(result['output'].splitlines()[:10])}
            else:
                return {"status": "error", "error_message": result['error_output']}
        except Exception as e:
            error_message = f"Error clearing temporary files: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}
    
    @staticmethod
    @error_handler
    def adjust_system_volume(volume_level):
        logging.info(f"Adjusting system volume to {volume_level}%")
        try:
            if platform.system() == 'Windows':
                devices = AudioUtilities.GetSpeakers()
                interface = devices.Activate(IAudioEndpointVolume._iid_, CLSCTX_ALL, None)
                volume = interface.QueryInterface(IAudioEndpointVolume)
                volume.SetMasterVolumeLevelScalar(volume_level / 100, None)
            elif platform.system() == 'Linux':
                os.system(f"amixer set Master {volume_level}%")
            else:
                return {"status": "error", "message": "Volume adjustment not supported on this platform"}
            return {"status": "success", "message": f"System volume adjusted to {volume_level}%"}
        except Exception as e:
            error_message = f"Error adjusting system volume: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}

    @staticmethod
    @error_handler
    def adjust_screen_brightness(brightness_level):
        logging.info(f"Adjusting screen brightness to {brightness_level}%")
        if platform.system() == 'Windows':
            powershell_path = os.getenv('POWERSHELL_EXE_PATH', r'C:\Windows\System32\WindowsPowerShell\v1.0\powershell.exe')
            command = (f'"{powershell_path}" (Get-WmiObject -Namespace root/WMI -Class WmiMonitorBrightnessMethods).WmiSetBrightness(1,{brightness_level})"')
            CommandExecutionAgent.execute_cmd(command)
        elif platform.system() == 'Linux':
            os.system(f"echo {brightness_level} | sudo tee /sys/class/backlight/*/brightness")
        else:
            return {"status": "error", "message": "Brightness adjustment not supported on this platform"}
        return {"status": "success", "message": f"Screen brightness adjusted to {brightness_level}%"}

    @staticmethod
    @error_handler
    def mute_system_volume():
        logging.info("Muting system volume")
        try:
            devices = AudioUtilities.GetSpeakers()
            interface = devices.Activate(IAudioEndpointVolume._iid_, CLSCTX_ALL, None)
            volume = cast(interface, POINTER(IAudioEndpointVolume))
            volume.SetMute(1, None)
            return {"status": "success", "message": "System volume muted"}
        except Exception as e:
            error_message = f"Error muting system volume: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}

    @staticmethod
    @error_handler
    def unmute_system_volume():
        logging.info("Unmuting system volume")
        try:
            devices = AudioUtilities.GetSpeakers()
            interface = devices.Activate(IAudioEndpointVolume._iid_, CLSCTX_ALL, None)
            volume = cast(interface, POINTER(IAudioEndpointVolume))
            volume.SetMute(0, None)
            return {"status": "success", "message": "System volume unmuted"}
        except Exception as e:
            error_message = f"Error unmuting system volume: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}
        
class MediaProcessingAgent:
    @staticmethod
    @error_handler
    def resize_image(input_output_pairs, width, height, maintain_aspect_ratio=False):
        if isinstance(input_output_pairs[0], str):
            input_output_pairs = [input_output_pairs]
        results = []
        for input_path, output_path in input_output_pairs:
            logging.info(f"Resizing image {input_path} to {width}x{height}")
            try:
                img = cv2.imread(input_path)
                if img is None:
                    raise ValueError("Image not found or invalid format")
                if maintain_aspect_ratio:
                    aspect_ratio = img.shape[1] / img.shape[0]
                    height = int(width / aspect_ratio)
                resized = cv2.resize(img, (width, height))
                cv2.imwrite(output_path, resized)
                success_message = f"Image resized and saved to {output_path}"
                logging.info(success_message)
                results.append({"status": "success", "input_path": input_path, "output_path": output_path, "message": success_message})
            except Exception as e:
                error_message = f"Error resizing image: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "input_path": input_path, "output_path": output_path, "error_message": error_message})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def apply_image_filter(input_output_filter_triples):
        if isinstance(input_output_filter_triples[0], str):
            input_output_filter_triples = [input_output_filter_triples]
        results = []
        for input_file, filter_name, output_file in input_output_filter_triples:
            logging.info(f"Applying filter {filter_name} to image {input_file}")
            try:
                with Image.open(input_file) as img:
                    filters = {
                        "BLUR": ImageFilter.BLUR,
                        "CONTOUR": ImageFilter.CONTOUR,
                        "EMBOSS": ImageFilter.EMBOSS,
                        "SHARPEN": ImageFilter.SHARPEN,
                        "SMOOTH": ImageFilter.SMOOTH,
                        "DETAIL": ImageFilter.DETAIL
                    }
                    if filter_name not in filters:
                        error_message = "Invalid filter name"
                        logging.error(error_message)
                        results.append({"status": "error", "input_file": input_file, "filter_name": filter_name, "output_file": output_file, "error_message": error_message})
                        continue
                    filtered_img = img.filter(filters[filter_name])
                    filtered_img.save(output_file)
                success_message = f"Image filtered and saved to {output_file}"
                logging.info(success_message)
                results.append({"status": "success", "input_file": input_file, "filter_name": filter_name, "output_file": output_file, "message": success_message})
            except Exception as e:
                error_message = f"Error applying filter to image: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "input_file": input_file, "filter_name": filter_name, "output_file": output_file, "error_message": error_message})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def convert_audio(input_output_format_triples):
        if isinstance(input_output_format_triples[0], str):
            input_output_format_triples = [input_output_format_triples]
        results = []
        for input_file, output_file, output_format in input_output_format_triples:
            logging.info(f"Converting audio {input_file} to {output_format} format")
            try:
                data, samplerate = sf.read(input_file)
                sf.write(output_file, data, samplerate, format=output_format)
                success_message = f"Audio converted to {output_format} and saved to {output_file}"
                logging.info(success_message)
                results.append({"status": "success", "input_file": input_file, "output_file": output_file, "output_format": output_format, "message": success_message})
            except Exception as e:
                error_message = f"Error converting audio: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "input_file": input_file, "output_file": output_file, "output_format": output_format, "error_message": error_message})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def trim_video(input_output_time_quads):
        if isinstance(input_output_time_quads[0], str):
            input_output_time_quads = [input_output_time_quads]
        results = []
        for input_file, output_file, start_time, end_time in input_output_time_quads:
            logging.info(f"Trimming video {input_file} from {start_time} to {end_time}")
            try:
                video = mp.VideoFileClip(input_file).subclip(start_time, end_time)
                video.write_videofile(output_file)
                success_message = f"Video trimmed and saved to {output_file}"
                logging.info(success_message)
                results.append({"status": "success", "input_file": input_file, "output_file": output_file, "start_time": start_time, "end_time": end_time, "message": success_message})
            except Exception as e:
                error_message = f"Error trimming video: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "input_file": input_file, "output_file": output_file, "start_time": start_time, "end_time": end_time, "error_message": error_message})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def extract_audio_from_video(input_output_pairs):
        if isinstance(input_output_pairs[0], str):
            input_output_pairs = [input_output_pairs]
        results = []
        for input_file, output_file in input_output_pairs:
            logging.info(f"Extracting audio from video {input_file}")
            try:
                video = mp.VideoFileClip(input_file)
                if video.audio is None:
                    raise ValueError("The video file does not contain any audio")
                audio = video.audio
                audio.write_audiofile(output_file)
                success_message = f"Audio extracted and saved to {output_file}"
                logging.info(success_message)
                results.append({"status": "success", "input_file": input_file, "output_file": output_file, "message": success_message})
            except Exception as e:
                error_message = f"Error extracting audio from video: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "input_file": input_file, "output_file": output_file, "error_message": error_message})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def merge_videos(input_files, output_file):
        logging.info(f"Merging videos {input_files} into {output_file}")
        try:
            clips = []
            base_clip = mp.VideoFileClip(input_files[0])
            base_size = base_clip.size

            for file in input_files:
                clip = mp.VideoFileClip(file)
                if clip.size != base_size:
                    clip = clip.resize(newsize=base_size)
                clips.append(clip)

            final_clip = mp.concatenate_videoclips(clips)
            final_clip.write_videofile(output_file)
            
            success_message = f"Videos merged and saved to {output_file}"
            logging.info(success_message)
            return {"status": "success", "input_files": input_files, "output_file": output_file, "message": success_message}
        except Exception as e:
            error_message = f"Error merging videos: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "input_files": input_files, "output_file": output_file, "error_message": error_message}

    @staticmethod
    @error_handler
    def add_audio_to_video(video_file, audio_file, output_file):
        logging.info(f"Adding audio {audio_file} to video {video_file}")
        try:
            video = mp.VideoFileClip(video_file)
            audio = mp.AudioFileClip(audio_file)
            final_video = video.set_audio(audio)
            final_video.write_videofile(output_file)
            success_message = f"Audio added to video and saved to {output_file}"
            logging.info(success_message)
            return {"status": "success", "message": success_message}
        except Exception as e:
            error_message = f"Error adding audio to video: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}
    
    @staticmethod
    @error_handler
    def video_subtitle(video_path, output_path, language='en', subtitle_file_path=None):
        logging.info(f"Adding subtitles to video: {video_path}")

        def read_srt_file(path):
            with open(path, 'r', encoding='utf-8') as file:
                content = file.read()
            return list(srt.parse(content))

        def generate_wrapped_text_clip(subtitle, video_width):
            max_text_width = int(video_width * 0.8)
            wrapper = textwrap.TextWrapper(width=70)
            wrapped_text = wrapper.fill(text=subtitle.content)

            return TextClip(txt=wrapped_text, fontsize=30, color='white', font='Arial',
                            align='center', method='caption', size=(max_text_width, None),
                            bg_color='rgba(0, 0, 0, 0.5)') \
                .set_position(('center', 'bottom')) \
                .set_duration((subtitle.end - subtitle.start).total_seconds()) \
                .set_start(subtitle.start.total_seconds())

        def timestamp_to_srt(seconds):
            return datetime.timedelta(seconds=int(seconds))

        try:
            video = VideoFileClip(video_path)
            if video.audio is None:
                raise ValueError("The video file does not contain any audio for generating subtitles")

            if subtitle_file_path is None:
                audio_path = video_path.replace(".mp4", ".wav")
                video.audio.write_audiofile(audio_path)

                model = WhisperModel("base")
                segments, _ = model.transcribe(audio_path, language=language)

                subtitles = []
                for i, segment in enumerate(segments):
                    start_time = timestamp_to_srt(segment.start)
                    end_time = timestamp_to_srt(segment.end)
                    subtitle = srt.Subtitle(index=i + 1, start=start_time, end=end_time, content=segment.text)
                    subtitles.append(subtitle)

                subtitle_file_path = video_path.replace(".mp4", ".srt")
                with open(subtitle_file_path, 'w', encoding='utf-8') as file:
                    file.write(srt.compose(subtitles))

                os.remove(audio_path)
            else:
                subtitles = read_srt_file(subtitle_file_path)

            txt_clips = [generate_wrapped_text_clip(subtitle, video.size[0]) for subtitle in subtitles]

            final_video = CompositeVideoClip([video, *txt_clips])
            final_video.write_videofile(output_path, codec='libx264', fps=video.fps)

            success_message = f"Subtitles added to video: {output_path}"
            logging.info(success_message)
            return {"status": "success", "output_path": output_path, "message": success_message}

        except Exception as e:
            error_message = f"Error adding subtitles to video: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}


class GUIAutomationAgent:
    @staticmethod
    @error_handler
    def mouse_click(x, y, clicks=1, interval=0.0, button='left', double_click=False):
        action = 'double click' if double_click else 'click'
        logging.info(f"Mouse {action} at ({x}, {y}) with {button} button")
        try:
            if double_click:
                pyautogui.doubleClick(x, y, interval=interval, button=button)
            else:
                pyautogui.click(x, y, clicks=clicks, interval=interval, button=button)
            success_message = f"Mouse {action} at ({x}, {y}) with {button} button"
            logging.info(success_message)
            return {"status": "success", "message": success_message}
        except Exception as e:
            error_message = f"Error during mouse {action}: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}

    @staticmethod
    @error_handler
    def type_text(text, interval=0.0, press_enter=False):
        logging.info(f"Typing text: {text}")
        try:
            pyautogui.typewrite(text, interval=interval)
            if press_enter:
                pyautogui.press('enter')
            success_message = f"Text typed: {text}"
            logging.info(success_message)
            return {"status": "success", "message": success_message}
        except Exception as e:
            error_message = f"Error typing text: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}

    @staticmethod
    @error_handler
    def press_key(key, combination=None):
        action = f"Pressing key: {key}" if not combination else f"Pressing keys: {combination + [key]}"
        logging.info(action)
        try:
            if combination:
                pyautogui.hotkey(*combination, key)
            else:
                pyautogui.press(key)
            success_message = action
            logging.info(success_message)
            return {"status": "success", "message": success_message}
        except Exception as e:
            error_message = f"Error pressing key: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}

    @staticmethod
    @error_handler
    def take_screenshot(output_path, region=None, active_window=False):
        logging.info(f"Taking screenshot. Region: {region}, Active window: {active_window}")
        try:
            if active_window:
                screenshot = ImageGrab.grab(bbox=pyautogui.getActiveWindow().box)
            else:
                screenshot = pyautogui.screenshot(region=region)
            screenshot.save(output_path)
            success_message = f"Screenshot saved to {output_path}"
            logging.info(success_message)
            return {"status": "success", "message": success_message}
        except Exception as e:
            error_message = f"Error taking screenshot: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}

    @staticmethod
    @error_handler
    def perform_ocr(image_paths, languages=['en']):
        if isinstance(image_paths, str):
            image_paths = [image_paths]
        results = []
        for image_path in image_paths:
            logging.info(f"Performing OCR on image {image_path} with languages {languages}")
            try:
                reader = easyocr.Reader(languages)
                result = reader.readtext(image_path)
                extracted_data = []
                for (bbox, text, _) in result:
                    (top_left, top_right, bottom_right, bottom_left) = bbox
                    middle_x = (top_left[0] + bottom_right[0]) / 2
                    middle_y = (top_left[1] + bottom_right[1]) / 2
                    extracted_data.append({"text": text, "middle_xy": (middle_x, middle_y)})
                
                success_message = f"OCR result for {image_path}: {extracted_data}"
                logging.info(success_message)
                results.append({"status": "success", "image_path": image_path, "data": extracted_data, "message": success_message})
            except Exception as e:
                error_message = f"Error performing OCR on {image_path}: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "image_path": image_path, "error_message": error_message})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def get_pixel_color(coordinates):
        if isinstance(coordinates[0], int):
            coordinates = [coordinates]
        results = []
        for x, y in coordinates:
            logging.info(f"Getting pixel color at ({x}, {y})")
            try:
                color = pyautogui.pixel(x, y)
                success_message = f"Pixel color at ({x}, {y}): {color}"
                logging.info(success_message)
                results.append({"status": "success", "coordinates": (x, y), "color": color})
            except Exception as e:
                error_message = f"Error getting pixel color at ({x}, {y}): {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "coordinates": (x, y), "error_message": error_message})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def move_mouse(x, y):
        logging.info(f"Moving mouse to ({x}, {y})")
        try:
            pyautogui.moveTo(x, y)
            success_message = f"Mouse moved to ({x}, {y})"
            logging.info(success_message)
            return {"status": "success", "message": success_message}
        except Exception as e:
            error_message = f"Error moving mouse: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}

    @staticmethod
    @error_handler
    def scroll(clicks, x=None, y=None):
        logging.info(f"Scrolling {clicks} clicks at ({x}, {y})")
        try:
            pyautogui.scroll(clicks, x, y)
            success_message = f"Scrolled {clicks} clicks at ({x}, {y})"
            logging.info(success_message)
            return {"status": "success", "message": success_message}
        except Exception as e:
            error_message = f"Error scrolling: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}

    @staticmethod
    @error_handler
    def hotkey(*args):
        logging.info(f"Pressing hotkey combination: {args}")
        try:
            pyautogui.hotkey(*args)
            success_message = f"Hotkey combination {args} pressed"
            logging.info(success_message)
            return {"status": "success", "message": success_message}
        except Exception as e:
            error_message = f"Error pressing hotkey combination: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}

    @staticmethod
    @error_handler
    def compare_images(image1_path, image2_path, diff_image_path=None):
        logging.info(f"Comparing images {image1_path} and {image2_path}")
        try:
            img1 = Image.open(image1_path).convert('RGB')
            img2 = Image.open(image2_path).convert('RGB')

            logging.info(f"Image 1 size: {img1.size}")
            logging.info(f"Image 2 size: {img2.size}")

            if img1.size != img2.size:
                img2 = img2.resize(img1.size, Image.LANCZOS)
                logging.info(f"Resized Image 2 to: {img1.size}")

            arr1 = np.array(img1)
            arr2 = np.array(img2)

            diff = np.abs(arr1 - arr2)
            diff_score = np.sum(diff)

            logging.info(f"Difference score: {diff_score}")

            diff_image = Image.fromarray(diff.astype('uint8'))
            if diff_image_path:
                diff_image.save(diff_image_path)
                logging.info(f"Difference image saved to: {diff_image_path}")

            success_message = f"Images compared. Difference score: {diff_score}"
            logging.info(success_message)
            return {"status": "success", "difference_score": int(diff_score), "message": success_message}
        except Exception as e:
            error_message = f"Error comparing images: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}

    @staticmethod
    @error_handler
    def find_image_on_screen(target_image_paths):
        if isinstance(target_image_paths, str):
            target_image_paths = [target_image_paths]
        results = []
        for target_image_path in target_image_paths:
            logging.info(f"Finding image {target_image_path} on screen")
            try:
                location = pyautogui.locateOnScreen(target_image_path)
                if location:
                    success_message = f"Image {target_image_path} found on screen at {location}"
                    logging.info(success_message)
                    results.append({"status": "success", "image_path": target_image_path, "location": location, "message": success_message})
                else:
                    error_message = f"Image {target_image_path} not found on screen"
                    logging.warning(error_message)
                    results.append({"status": "error", "image_path": target_image_path, "error_message": error_message})
            except Exception as e:
                error_message = f"Error finding image {target_image_path} on screen: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "image_path": target_image_path, "error_message": error_message})
        return results if len(results) > 1 else results[0]
    
    @staticmethod
    @error_handler
    def click_image_on_screen(target_image_path):
        logging.info(f"Clicking on image {target_image_path} on screen")
        try:
            location = pyautogui.locateCenterOnScreen(target_image_path)
            if location:
                pyautogui.click(location)
                success_message = f"Clicked on image at {location}"
                logging.info(success_message)
                return {"status": "success", "location": location, "message": success_message}
            else:
                error_message = "Image not found on screen"
                logging.warning(error_message)
                return {"status": "error", "error_message": error_message}
        except Exception as e:
            error_message = f"Error clicking on image on screen: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}
    
    @staticmethod
    @error_handler
    def drag_mouse(start_x, start_y, end_x, end_y, duration=0.5):
        logging.info(f"Dragging mouse from ({start_x}, {start_y}) to ({end_x}, {end_y})")
        try:
            pyautogui.moveTo(start_x, start_y)
            pyautogui.dragTo(end_x, end_y, duration=duration, button='left')
            success_message = f"Mouse dragged from ({start_x}, {start_y}) to ({end_x}, {end_y})"
            logging.info(success_message)
            return {"status": "success", "message": success_message}
        except Exception as e:
            error_message = f"Error dragging mouse: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}

class ProcessManagementAgent:
    @staticmethod
    @error_handler
    def list_processes():
        logging.info("Listing all running processes")
        try:
            processes = [{'pid': p.pid, 'name': p.name(), 'cpu_percent': p.cpu_percent(), 'memory_percent': p.memory_percent()} for p in psutil.process_iter(['pid', 'name'])]
            success_message = "Processes listed successfully"
            logging.info(success_message)
            return {"status": "success", "processes": processes, "message": success_message}
        except Exception as e:
            error_message = f"Error listing processes: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}

    @staticmethod
    @error_handler
    def start_process(commands, paths=None, threaded=False):
        if isinstance(commands, str):
            commands = [commands]
        if paths is None:
            paths = [None] * len(commands)
        elif isinstance(paths, str):
            paths = [paths]
        
        results = []
        for command, path in zip(commands, paths):
            logging.info(f"Starting process with command: {command}, Path: {path}, Threaded: {threaded}")
            process_info = {"status": "error", "message": "Failed to start process"}

            def run_process():
                nonlocal process_info
                try:
                    if path:
                        full_command = f"cd {path} && {command}"
                    else:
                        full_command = command
                    process = subprocess.Popen(full_command, shell=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
                    process_info = {
                        "status": "success",
                        "pid": process.pid,
                        "stdout": process.stdout,
                        "stderr": process.stderr,
                        "message": f"Process started with PID: {process.pid}"
                    }
                    logging.info(process_info["message"])
                except Exception as e:
                    process_info = {"status": "error", "message": f"Error starting process: {str(e)}"}
                    logging.error(process_info["message"])

            if threaded:
                thread = threading.Thread(target=run_process)
                thread.start()
                thread.join()
            else:
                run_process()
            results.append(process_info)
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def stop_process(pids, force=False):
        if isinstance(pids, int):
            pids = [pids]
        results = []
        for pid in pids:
            logging.info(f"Stopping process with PID: {pid}, Force: {force}")
            try:
                process = psutil.Process(pid)
                if force:
                    process.kill()
                else:
                    process.terminate()
                process.wait()
                success_message = f"Process with PID {pid} terminated"
                logging.info(success_message)
                results.append({"status": "success", "pid": pid, "message": success_message})
            except psutil.NoSuchProcess:
                error_message = f"No process found with PID {pid}"
                logging.error(error_message)
                results.append({"status": "error", "pid": pid, "error_message": error_message})
            except Exception as e:
                error_message = f"Error stopping process with PID {pid}: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "pid": pid, "error_message": error_message})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def get_process_info(pids):
        if isinstance(pids, int):
            pids = [pids]
        results = []
        for pid in pids:
            logging.info(f"Getting information for process with PID: {pid}")
            try:
                process = psutil.Process(pid)
                process_info = {
                    'name': process.name(),
                    'status': process.status(),
                    'cpu_percent': process.cpu_percent(),
                    'memory_percent': process.memory_percent(),
                    'create_time': process.create_time(),
                    'open_files': process.open_files(),
                    'connections': process.connections()
                }
                success_message = f"Process information retrieved for PID {pid}"
                logging.info(success_message)
                results.append({"status": "success", "pid": pid, "process_info": process_info, "message": success_message})
            except psutil.NoSuchProcess:
                error_message = f"No process found with PID {pid}"
                logging.error(error_message)
                results.append({"status": "error", "pid": pid, "error_message": error_message})
            except Exception as e:
                error_message = f"Error getting process information for PID {pid}: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "pid": pid, "error_message": error_message})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def restart_process(pids, threaded=False):
        if isinstance(pids, int):
            pids = [pids]
        results = []
        for pid in pids:
            logging.info(f"Restarting process with PID: {pid}, Threaded: {threaded}")
            def restart():
                try:
                    process = psutil.Process(pid)
                    command = process.cmdline()
                    process.terminate()
                    process.wait()
                    new_process = subprocess.Popen(command, shell=True)
                    success_message = f"Process with PID {pid} restarted with new PID {new_process.pid}"
                    logging.info(success_message)
                    return {"status": "success", "old_pid": pid, "new_pid": new_process.pid, "message": success_message}
                except psutil.NoSuchProcess:
                    error_message = f"No process found with PID {pid}"
                    logging.error(error_message)
                    return {"status": "error", "pid": pid, "error_message": error_message}
                except Exception as e:
                    error_message = f"Error restarting process with PID {pid}: {str(e)}"
                    logging.error(error_message)
                    return {"status": "error", "pid": pid, "error_message": error_message}
            
            if threaded:
                thread = threading.Thread(target=restart)
                thread.start()
                results.append({"status": "success", "pid": pid, "message": f"Process restart initiated in a separate thread for PID {pid}"})
            else:
                results.append(restart())
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def monitor_process(pids, interval=1, threaded=True):
        if isinstance(pids, int):
            pids = [pids]
        results = []
        for pid in pids:
            logging.info(f"Monitoring process with PID: {pid} every {interval} seconds, Threaded: {threaded}")
            def monitor():
                try:
                    process = psutil.Process(pid)
                    while process.is_running():
                        info = {
                            'cpu_percent': process.cpu_percent(interval=interval),
                            'memory_percent': process.memory_percent(),
                            'status': process.status()
                        }
                        logging.info(f"Process {pid} info: {info}")
                        time.sleep(interval)
                except psutil.NoSuchProcess:
                    logging.warning(f"Process with PID {pid} no longer exists")
                except Exception as e:
                    logging.error(f"Error monitoring process with PID {pid}: {str(e)}")

            if threaded:
                thread = threading.Thread(target=monitor)
                thread.start()
                results.append({"status": "success", "pid": pid, "message": f"Started monitoring process with PID: {pid} in a separate thread"})
            else:
                monitor()
                results.append({"status": "success", "pid": pid, "message": f"Started monitoring process with PID: {pid}"})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def set_resource_alert(pids, cpu_threshold=80, memory_threshold=80, check_interval=5, threaded=True):
        if isinstance(pids, int):
            pids = [pids]
        results = []
        for pid in pids:
            logging.info(f"Setting resource alerts for process {pid}. CPU threshold: {cpu_threshold}%, Memory threshold: {memory_threshold}%")
            def monitor():
                try:
                    process = psutil.Process(pid)
                    while process.is_running():
                        cpu_usage = process.cpu_percent(interval=check_interval)
                        memory_usage = process.memory_percent()
                        if cpu_usage > cpu_threshold:
                            logging.warning(f"Process {pid} CPU usage high: {cpu_usage}%")
                        if memory_usage > memory_threshold:
                            logging.warning(f"Process {pid} memory usage high: {memory_usage}%")
                        time.sleep(check_interval)
                except psutil.NoSuchProcess:
                    logging.warning(f"Process with PID {pid} no longer exists")
                except Exception as e:
                    logging.error(f"Error monitoring resource usage for process with PID {pid}: {str(e)}")

            if threaded:
                thread = threading.Thread(target=monitor)
                thread.start()
                results.append({"status": "success", "pid": pid, "message": f"Started resource monitoring for process with PID: {pid} in a separate thread"})
            else:
                monitor()
                results.append({"status": "success", "pid": pid, "message": f"Started resource monitoring for process with PID: {pid}"})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def log_process_output(commands, log_files, threaded=True):
        if isinstance(commands, str):
            commands = [commands]
        if isinstance(log_files, str):
            log_files = [log_files]
        results = []
        for command, log_file in zip(commands, log_files):
            logging.info(f"Starting process with command: {command} and logging output to {log_file}, Threaded: {threaded}")
            def log_output():
                try:
                    os.makedirs(os.path.dirname(log_file), exist_ok=True)
                    with open(log_file, 'w') as f:
                        process = subprocess.Popen(command, shell=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
                        for line in iter(process.stdout.readline, ''):
                            f.write(line)
                            logging.info(line.strip())
                        process.stdout.close()
                        process.wait()
                    success_message = f"Process output logged to {log_file}"
                    logging.info(success_message)
                    return {"status": "success", "command": command, "log_file": log_file, "message": success_message}
                except Exception as e:
                    error_message = f"Error logging process output for command {command}: {str(e)}"
                    logging.error(error_message)
                    return {"status": "error", "command": command, "log_file": log_file, "error_message": error_message}
            
            if threaded:
                thread = threading.Thread(target=log_output)
                thread.start()
                results.append({"status": "success", "command": command, "log_file": log_file, "message": f"Started logging output for command: {command} in a separate thread"})
            else:
                results.append(log_output())
        return results if len(results) > 1 else results[0]

        
class NetworkOperationsAgent:
    nmap_path = os.getenv('NMAP_PATH', r'C:\Program Files (x86)\Nmap\nmap.exe')

    @staticmethod
    @error_handler
    def get_ip_address():
        logging.info("Retrieving local IP address")
        hostname = socket.gethostname()
        local_ip = socket.gethostbyname(hostname)
        success_message = f"Local IP address retrieved: {local_ip}"
        logging.info(success_message)
        return {"status": "success", "local_ip": local_ip, "hostname": hostname, "message": success_message}

    @staticmethod
    @error_handler
    def ping(hosts, count=4, timeout=5):
        if isinstance(hosts, str):
            hosts = [hosts]
        results = []
        for host in hosts:
            logging.info(f"Pinging host: {host} with count: {count}")
            param = '-n' if platform.system().lower() == 'windows' else '-c'
            command = ['ping', param, str(count), host]
            try:
                result = subprocess.run(command, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=timeout)
                try:
                    output = result.stdout.decode('utf-8')
                except UnicodeDecodeError:
                    output = result.stdout.decode('latin-1')
                success_message = f"Ping result for host {host}: {output}"
                logging.info(success_message)
                results.append({"status": "success", "host": host, "output": output, "message": success_message})
            except subprocess.TimeoutExpired:
                error_message = f"Ping to {host} timed out after {timeout} seconds"
                logging.error(error_message)
                results.append({"status": "error", "host": host, "error_message": error_message})
            except Exception as e:
                error_message = f"Error pinging {host}: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "host": host, "error_message": error_message})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def port_scan(targets, port_range="22-80"):
        if isinstance(targets, str):
            targets = [targets]
        results = []
        for target in targets:
            logging.info(f"Scanning ports on target: {target}, Port range: {port_range}")
            start_port, end_port = map(int, port_range.split('-'))
            open_ports = []

            def check_port(port):
                with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as sock:
                    sock.settimeout(0.1)
                    result = sock.connect_ex((target, port))
                    if result == 0:
                        return port
                return None

            with concurrent.futures.ThreadPoolExecutor(max_workers=100) as executor:
                futures = [executor.submit(check_port, port) for port in range(start_port, end_port + 1)]
                for future in concurrent.futures.as_completed(futures):
                    port = future.result()
                    if port:
                        open_ports.append(port)

            success_message = f"Open ports on {target}: {open_ports}"
            logging.info(success_message)
            results.append({"status": "success", "target": target, "open_ports": open_ports, "message": success_message})
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def get_public_ip():
        logging.info("Retrieving public IP address")
        try:
            public_ip = requests.get('https://api.ipify.org', timeout=5).text
            isp_info = requests.get(f'https://ipinfo.io/{public_ip}/json', timeout=5).json()
            success_message = f"Public IP address retrieved: {public_ip}, ISP: {isp_info.get('org', 'N/A')}"
            logging.info(success_message)
            return {"status": "success", "public_ip": public_ip, "isp_info": isp_info, "message": success_message}
        except requests.RequestException as e:
            error_message = f"Error retrieving public IP: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}

    @staticmethod
    @error_handler
    def check_internet_connection(test_urls=['https://www.google.com', 'https://www.cloudflare.com', 'https://www.amazon.com']):
        logging.info(f"Checking internet connection")
        for url in test_urls:
            try:
                response = requests.get(url, timeout=5)
                latency = response.elapsed.total_seconds()
                success_message = f"Internet connection is up. Latency to {url}: {latency} seconds"
                logging.info(success_message)
                return {"status": "success", "test_url": url, "latency": latency, "message": success_message}
            except requests.RequestException:
                continue
        error_message = "Internet connection is down or unreliable"
        logging.error(error_message)
        return {"status": "error", "error_message": error_message}

    @staticmethod
    @error_handler
    def perform_speedtest(timeout=20):
        logging.info(f"Performing internet speed test with timeout: {timeout} seconds")
        try:
            st = speedtest.Speedtest()
            st.get_best_server()
            
            def run_test():
                download_speed = st.download() / 1_000_000 
                upload_speed = st.upload() / 1_000_000
                ping = st.results.ping
                return download_speed, upload_speed, ping

            with concurrent.futures.ThreadPoolExecutor() as executor:
                future = executor.submit(run_test)
                download_speed, upload_speed, ping = future.result(timeout=timeout)

            success_message = "Speed test completed"
            logging.info(success_message)
            return {
                "status": "success",
                "download_speed": round(download_speed, 2),
                "upload_speed": round(upload_speed, 2),
                "ping": round(ping, 2),
                "message": success_message
            }
        except concurrent.futures.TimeoutError:
            error_message = f"Speed test timed out after {timeout} seconds"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}
        except Exception as e:
            error_message = f"Error performing speed test: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "error_message": error_message}

class SpeechAgent:
    def __init__(self):
        self.engine = pyttsx3.init()
        self.model = None
        self.voices = self.engine.getProperty('voices')

    @staticmethod
    def is_cuda_available():
        try:
            import torch
            return torch.cuda.is_available()
        except ImportError:
            return False

    def load_whisper_model(self):
        if self.model is None:
            self.model = WhisperModel("base", device="cuda" if self.is_cuda_available() else "cpu")

    def text_to_speech(self, text: str, language: str = 'en', voice_id: Optional[str] = None, 
                       rate: Optional[int] = None, volume: Optional[float] = None, 
                       output_file: Optional[str] = None):
        if voice_id:
            self.engine.setProperty('voice', voice_id)
        else:
            language_voice = next((voice for voice in self.voices if language.lower() in voice.id.lower()), None)
            if language_voice:
                self.engine.setProperty('voice', language_voice.id)
            else:
                return {"status": "error", "message": f"No voice found for language: {language}"}

        if rate is not None:
            self.engine.setProperty('rate', rate)
        if volume is not None:
            self.engine.setProperty('volume', volume)
        
        if output_file:
            output_dir = os.path.dirname(output_file)
            if not os.path.exists(output_dir):
                os.makedirs(output_dir)
            self.engine.save_to_file(text, output_file)
        else:
            self.engine.say(text)
        
        self.engine.runAndWait()
        return {"status": "success", "message": "Text to speech conversion successful"}

    @staticmethod
    def speech_to_text(audio_file: Optional[str] = None, language: str = 'en-US', 
                       timeout: int = 5, phrase_time_limit: Optional[int] = None):
        recognizer = sr.Recognizer()
        
        if audio_file:
            if not os.path.exists(audio_file):
                return {"status": "error", "message": f"Audio file not found: {audio_file}"}
            with sr.AudioFile(audio_file) as source:
                audio = recognizer.record(source)
        else:
            try:
                with sr.Microphone() as source:
                    print("Adjusting for ambient noise. Please wait...")
                    recognizer.adjust_for_ambient_noise(source, duration=1)
                    print("Listening...")
                    audio = recognizer.listen(source, timeout=timeout, phrase_time_limit=phrase_time_limit)
            except sr.WaitTimeoutError:
                return {"status": "error", "message": "Listening timed out. No speech detected."}
            except Exception as e:
                return {"status": "error", "message": f"Error accessing the microphone: {str(e)}"}
        
        try:
            text = recognizer.recognize_google(audio, language=language)
            return {"status": "success", "text": text, "message": "Speech to text conversion successful"}
        except sr.UnknownValueError:
            return {"status": "error", "message": "Speech recognition could not understand audio"}
        except sr.RequestError:
            return {"status": "error", "message": "Could not request results from speech recognition service"}

    def whisper_transcription(self, audio_files: Optional[Union[str, List[str]]] = None, duration: int = 10, 
                              output_files: Optional[Union[str, List[str]]] = None, language: Optional[str] = None):
        import sounddevice as sd
        self.load_whisper_model()
        
        if isinstance(audio_files, str):
            audio_files = [audio_files]
        if isinstance(output_files, str):
            output_files = [output_files]
        
        results = []
        
        for i, audio_file in enumerate(audio_files or [None]):
            if audio_file:
                segments, info = self.model.transcribe(audio_file, language=language)
            else:
                print(f"Starting Whisper live transcription for {duration} seconds...")
                audio_data = []
                
                def audio_callback(indata, frames, time, status):
                    if status:
                        print(status, file=sys.stderr)
                    audio_data.extend(indata[:, 0])
                
                with sd.InputStream(callback=audio_callback, channels=1, samplerate=16000):
                    sd.sleep(int(duration * 1000))
                
                audio_array = np.array(audio_data)
                segments, info = self.model.transcribe(audio_array, language=language)
            
            text = " ".join([segment.text for segment in segments])
            print(f"Transcription: {text}")
            
            if output_files and i < len(output_files):
                with open(output_files[i], 'w') as f:
                    f.write(text)
                print(f"Transcription saved to {output_files[i]}")
            
            results.append({"status": "success", "transcription": text, "message": "Transcription completed"})
        
        return results if len(results) > 1 else results[0]

    def transcribe_long_audio(self, audio_files: Union[str, List[str]], chunk_duration: int = 30, language: Optional[str] = None):
        self.load_whisper_model()
        
        if isinstance(audio_files, str):
            audio_files = [audio_files]
        
        results = []
        
        for audio_file in audio_files:
            full_transcription = []
            with sr.AudioFile(audio_file) as source:
                audio_duration = source.DURATION
                for i in range(0, int(audio_duration), chunk_duration):
                    chunk = sr.AudioFile(audio_file)
                    with chunk as source:
                        audio = sr.Recognizer().record(source, offset=i, duration=min(chunk_duration, audio_duration - i))
                        audio_data = np.frombuffer(audio.get_raw_data(), dtype=np.int16).flatten().astype(np.float32) / 32768.0
                        segments, _ = self.model.transcribe(audio_data, language=language)
                        chunk_text = " ".join([segment.text for segment in segments])
                        full_transcription.append(chunk_text)
            
            full_text = " ".join(full_transcription)
            results.append({"status": "success", "transcription": full_text, "message": "Long audio transcription completed"})
        
        return results if len(results) > 1 else results[0]

    def list_available_voices(self):
        voices = self.engine.getProperty('voices')
        available_voices = [{"id": voice.id, "name": voice.name, "languages": voice.languages} for voice in voices]
        return {"status": "success", "voices": available_voices, "message": "Available voices listed"}
    
class DatabaseOperationsAgent:
    @staticmethod
    @error_handler
    def execute_query(db_files, queries, params=None):
        if isinstance(db_files, str):
            db_files = [db_files]
        if isinstance(queries, str):
            queries = [queries]
        if params is None:
            params = [None] * len(queries)
        elif not isinstance(params[0], (list, tuple)):
            params = [params]
        
        results = []
        for db_file, query, param in zip(db_files, queries, params):
            logging.info(f"Executing query: {query} on {db_file}")
            try:
                conn = sqlite3.connect(db_file)
                cursor = conn.cursor()
                if param:
                    cursor.execute(query, param)
                else:
                    cursor.execute(query)
                result = cursor.fetchall()
                conn.commit()
                conn.close()
                success_message = f"Query executed successfully: {query}"
                logging.info(success_message)
                results.append({"status": "success", "result": result, "message": success_message})
            except sqlite3.Error as e:
                error_message = f"SQLite error executing query: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "error_message": error_message})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def create_table(db_files, table_names, columns):
        if isinstance(db_files, str):
            db_files = [db_files]
        if isinstance(table_names, str):
            table_names = [table_names]
        if isinstance(columns[0], str):
            columns = [columns] * len(db_files)
        
        results = []
        for db_file, table_name, cols in zip(db_files, table_names, columns):
            logging.info(f"Creating table: {table_name} with columns: {cols} in {db_file}")
            query = f"CREATE TABLE IF NOT EXISTS {table_name} ({', '.join(cols)})"
            results.append(DatabaseOperationsAgent.execute_query(db_file, query))
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def insert_data(db_files, table_names, data):
        if isinstance(db_files, str):
            db_files = [db_files]
        if isinstance(table_names, str):
            table_names = [table_names]
        if not isinstance(data[0], (list, tuple)):
            data = [data]
        
        results = []
        for db_file, table_name, row in zip(db_files, table_names, data):
            logging.info(f"Inserting data into {table_name} in {db_file}: {row}")
            placeholders = ', '.join(['?' for _ in row])
            query = f"INSERT INTO {table_name} VALUES ({placeholders})"
            results.append(DatabaseOperationsAgent.execute_query(db_file, query, row))
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def bulk_insert_data(db_files, table_names, data_lists):
        if isinstance(db_files, str):
            db_files = [db_files]
        if isinstance(table_names, str):
            table_names = [table_names]
        if not isinstance(data_lists[0][0], (list, tuple)):
            data_lists = [data_lists]
        
        results = []
        for db_file, table_name, data_list in zip(db_files, table_names, data_lists):
            logging.info(f"Bulk inserting data into {table_name} in {db_file}")
            placeholders = ', '.join(['?' for _ in data_list[0]])
            query = f"INSERT INTO {table_name} VALUES ({placeholders})"
            try:
                conn = sqlite3.connect(db_file)
                cursor = conn.cursor()
                cursor.executemany(query, data_list)
                conn.commit()
                conn.close()
                success_message = "Bulk insert executed successfully"
                logging.info(success_message)
                results.append({"status": "success", "message": success_message})
            except sqlite3.Error as e:
                error_message = f"SQLite error executing bulk insert: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "error_message": error_message})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def retrieve_data(db_files, table_names, columns='*', conditions=None):
        if isinstance(db_files, str):
            db_files = [db_files]
        if isinstance(table_names, str):
            table_names = [table_names]
        if isinstance(columns, str):
            columns = [columns] * len(db_files)
        if conditions is None:
            conditions = [None] * len(db_files)
        elif isinstance(conditions, str):
            conditions = [conditions] * len(db_files)
        
        results = []
        for db_file, table_name, cols, cond in zip(db_files, table_names, columns, conditions):
            logging.info(f"Retrieving data from {table_name} in {db_file} with conditions: {cond}")
            query = f"SELECT {cols} FROM {table_name}"
            if cond:
                query += f" WHERE {cond}"
            results.append(DatabaseOperationsAgent.execute_query(db_file, query))
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def update_data(db_files, table_names, updates, conditions=None):
        if isinstance(db_files, str):
            db_files = [db_files]
        if isinstance(table_names, str):
            table_names = [table_names]
        if isinstance(updates, str):
            updates = [updates] * len(db_files)
        if conditions is None:
            conditions = [None] * len(db_files)
        elif isinstance(conditions, str):
            conditions = [conditions] * len(db_files)
        
        results = []
        for db_file, table_name, update, cond in zip(db_files, table_names, updates, conditions):
            logging.info(f"Updating data in {table_name} in {db_file} with updates: {update} and conditions: {cond}")
            query = f"UPDATE {table_name} SET {update}"
            if cond:
                query += f" WHERE {cond}"
            results.append(DatabaseOperationsAgent.execute_query(db_file, query))
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def delete_data(db_files, table_names, conditions):
        if isinstance(db_files, str):
            db_files = [db_files]
        if isinstance(table_names, str):
            table_names = [table_names]
        if isinstance(conditions, str):
            conditions = [conditions] * len(db_files)
        
        results = []
        for db_file, table_name, cond in zip(db_files, table_names, conditions):
            logging.info(f"Deleting data from {table_name} in {db_file} with conditions: {cond}")
            query = f"DELETE FROM {table_name} WHERE {cond}"
            results.append(DatabaseOperationsAgent.execute_query(db_file, query))
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def check_table_exists(db_files, table_names):
        if isinstance(db_files, str):
            db_files = [db_files]
        if isinstance(table_names, str):
            table_names = [table_names]
        
        results = []
        for db_file, table_name in zip(db_files, table_names):
            logging.info(f"Checking if table {table_name} exists in {db_file}")
            query = f"SELECT name FROM sqlite_master WHERE type='table' AND name='{table_name}'"
            result = DatabaseOperationsAgent.execute_query(db_file, query)
            exists = bool(result['result'])
            success_message = f"Table {table_name} exists: {exists}"
            logging.info(success_message)
            results.append({"status": "success", "exists": exists, "message": success_message})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def retrieve_schema(db_files, table_names):
        if isinstance(db_files, str):
            db_files = [db_files]
        if isinstance(table_names, str):
            table_names = [table_names]
        
        results = []
        for db_file, table_name in zip(db_files, table_names):
            logging.info(f"Retrieving schema of table {table_name} in {db_file}")
            query = f"PRAGMA table_info({table_name})"
            results.append(DatabaseOperationsAgent.execute_query(db_file, query))
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def backup_database(db_files, backup_files):
        if isinstance(db_files, str):
            db_files = [db_files]
        if isinstance(backup_files, str):
            backup_files = [backup_files]
        
        results = []
        for db_file, backup_file in zip(db_files, backup_files):
            logging.info(f"Backing up database {db_file} to {backup_file}")
            try:
                conn = sqlite3.connect(db_file)
                with open(backup_file, 'w') as f:
                    for line in conn.iterdump():
                        f.write(f"{line}\n")
                conn.close()
                success_message = f"Database backup successful to {backup_file}"
                logging.info(success_message)
                results.append({"status": "success", "message": success_message})
            except sqlite3.Error as e:
                error_message = f"SQLite error during backup: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "error_message": error_message})
        
        return results if len(results) > 1 else results[0]

class WindowManagementAgent:
    @staticmethod
    @error_handler
    def focus_window(window_titles):
        if isinstance(window_titles, str):
            window_titles = [window_titles]
        
        results = []
        for window_title in window_titles:
            logging.info(f"Focusing window: {window_title}")
            hwnd = win32gui.FindWindow(None, window_title)
            if not hwnd:
                error_message = f"Window not found: {window_title}"
                logging.error(error_message)
                results.append({"status": "error", "message": error_message})
                continue
            
            try:
                win32gui.ShowWindow(hwnd, win32con.SW_RESTORE)
                ctypes.windll.user32.SetForegroundWindow(hwnd)
                results.append({"status": "success", "message": f"Focused window: {window_title}"})
            except Exception as e:
                error_message = f"Failed to focus window: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "message": error_message})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def get_open_windows():
        logging.info("Retrieving list of open windows")
        def callback(hwnd, windows):
            if win32gui.IsWindowVisible(hwnd):
                windows.append(win32gui.GetWindowText(hwnd))
            return True
        windows = []
        win32gui.EnumWindows(callback, windows)
        logging.info(f"Open windows: {windows}")
        return {"status": "success", "windows": windows, "message": "Retrieved list of open windows"}

    @staticmethod
    @error_handler
    def minimize_window(window_titles):
        if isinstance(window_titles, str):
            window_titles = [window_titles]
        
        results = []
        for window_title in window_titles:
            logging.info(f"Minimizing window: {window_title}")
            hwnd = win32gui.FindWindow(None, window_title)
            if not hwnd:
                error_message = f"Window not found: {window_title}"
                logging.error(error_message)
                results.append({"status": "error", "message": error_message})
                continue
            win32gui.ShowWindow(hwnd, win32con.SW_MINIMIZE)
            results.append({"status": "success", "message": f"Minimized window: {window_title}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def maximize_window(window_titles):
        if isinstance(window_titles, str):
            window_titles = [window_titles]
        
        results = []
        for window_title in window_titles:
            logging.info(f"Maximizing window: {window_title}")
            hwnd = win32gui.FindWindow(None, window_title)
            if not hwnd:
                error_message = f"Window not found: {window_title}"
                logging.error(error_message)
                results.append({"status": "error", "message": error_message})
                continue
            win32gui.ShowWindow(hwnd, win32con.SW_MAXIMIZE)
            results.append({"status": "success", "message": f"Maximized window: {window_title}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def close_window(window_titles):
        if isinstance(window_titles, str):
            window_titles = [window_titles]
        
        results = []
        for window_title in window_titles:
            logging.info(f"Closing window: {window_title}")
            hwnd = win32gui.FindWindow(None, window_title)
            if not hwnd:
                error_message = f"Window not found: {window_title}"
                logging.error(error_message)
                results.append({"status": "error", "message": error_message})
                continue
            win32gui.PostMessage(hwnd, win32con.WM_CLOSE, 0, 0)
            results.append({"status": "success", "message": f"Closed window: {window_title}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def resize_move_window(window_titles, positions):
        if isinstance(window_titles, str):
            window_titles = [window_titles]
        if isinstance(positions[0], int):
            positions = [positions]
        
        results = []
        for window_title, (x, y, width, height) in zip(window_titles, positions):
            logging.info(f"Resizing and moving window: {window_title} to ({x}, {y}, {width}, {height})")
            hwnd = win32gui.FindWindow(None, window_title)
            if not hwnd:
                error_message = f"Window not found: {window_title}"
                logging.error(error_message)
                results.append({"status": "error", "message": error_message})
                continue
            win32gui.MoveWindow(hwnd, x, y, width, height, True)
            results.append({"status": "success", "message": f"Resized and moved window: {window_title}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def list_window_titles():
        logging.info("Listing titles of all open windows")
        def callback(hwnd, titles):
            if win32gui.IsWindowVisible(hwnd):
                titles.append(win32gui.GetWindowText(hwnd))
            return True
        titles = []
        win32gui.EnumWindows(callback, titles)
        logging.info(f"Window titles: {titles}")
        return {"status": "success", "titles": titles, "message": "Listed titles of all open windows"}

    @staticmethod
    @error_handler
    def bring_window_to_front(window_titles):
        if isinstance(window_titles, str):
            window_titles = [window_titles]
        
        results = []
        for window_title in window_titles:
            logging.info(f"Bringing window to front: {window_title}")
            hwnd = win32gui.FindWindow(None, window_title)
            if not hwnd:
                error_message = f"Window not found: {window_title}"
                logging.error(error_message)
                results.append({"status": "error", "message": error_message})
                continue
            win32gui.SetWindowPos(hwnd, win32con.HWND_TOPMOST, 0, 0, 0, 0, win32con.SWP_NOMOVE | win32con.SWP_NOSIZE)
            win32gui.SetWindowPos(hwnd, win32con.HWND_NOTOPMOST, 0, 0, 0, 0, win32con.SWP_SHOWWINDOW | win32con.SWP_NOMOVE | win32con.SWP_NOSIZE)
            results.append({"status": "success", "message": f"Brought window to front: {window_title}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def check_window_state(window_titles):
        if isinstance(window_titles, str):
            window_titles = [window_titles]
        
        results = []
        for window_title in window_titles:
            logging.info(f"Checking state of window: {window_title}")
            hwnd = win32gui.FindWindow(None, window_title)
            if not hwnd:
                error_message = f"Window not found: {window_title}"
                logging.error(error_message)
                results.append({"status": "error", "message": error_message})
                continue
            
            placement = win32gui.GetWindowPlacement(hwnd)
            if placement[1] == win32con.SW_SHOWMINIMIZED:
                state = "minimized"
            elif placement[1] == win32con.SW_SHOWMAXIMIZED:
                state = "maximized"
            else:
                state = "normal"
            
            success_message = f"Window {window_title} is {state}"
            logging.info(success_message)
            results.append({"status": "success", "state": state, "message": success_message})
        
        return results if len(results) > 1 else results[0]


    @staticmethod
    @error_handler
    def capture_window_screenshot(window_titles, save_paths):
        if isinstance(window_titles, str):
            window_titles = [window_titles]
        if isinstance(save_paths, str):
            save_paths = [save_paths]
        
        results = []
        for window_title, save_path in zip(window_titles, save_paths):
            logging.info(f"Capturing screenshot of window: {window_title}")
            hwnd = win32gui.FindWindow(None, window_title)
            if not hwnd:
                error_message = f"Window not found: {window_title}"
                logging.error(error_message)
                results.append({"status": "error", "window": window_title, "message": error_message})
                continue

            left, top, right, bottom = win32gui.GetWindowRect(hwnd)
            width = right - left
            height = bottom - top

            hwndDC = win32gui.GetWindowDC(hwnd)
            mfcDC = win32ui.CreateDCFromHandle(hwndDC)
            saveDC = mfcDC.CreateCompatibleDC()

            saveBitMap = win32ui.CreateBitmap()
            saveBitMap.CreateCompatibleBitmap(mfcDC, width, height)

            saveDC.SelectObject(saveBitMap)

            result = windll.user32.PrintWindow(hwnd, saveDC.GetSafeHdc(), 0)

            bmpinfo = saveBitMap.GetInfo()
            bmpstr = saveBitMap.GetBitmapBits(True)

            im = Image.frombuffer(
                'RGB',
                (bmpinfo['bmWidth'], bmpinfo['bmHeight']),
                bmpstr, 'raw', 'BGRX', 0, 1)

            win32gui.DeleteObject(saveBitMap.GetHandle())
            saveDC.DeleteDC()
            mfcDC.DeleteDC()
            win32gui.ReleaseDC(hwnd, hwndDC)

            im.save(save_path)

            if result == 1:
                success_message = f"Screenshot saved to {save_path}"
                logging.info(success_message)
                results.append({"status": "success", "window": window_title, "path": save_path, "message": success_message})
            else:
                error_message = f"Failed to capture screenshot for {window_title}"
                logging.error(error_message)
                results.append({"status": "error", "window": window_title, "message": error_message})

        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def change_window_title(window_titles, new_titles):
        if isinstance(window_titles, str):
            window_titles = [window_titles]
        if isinstance(new_titles, str):
            new_titles = [new_titles]
        
        results = []
        for window_title, new_title in zip(window_titles, new_titles):
            logging.info(f"Changing window title from {window_title} to {new_title}")
            hwnd = win32gui.FindWindow(None, window_title)
            if not hwnd:
                error_message = f"Window not found: {window_title}"
                logging.error(error_message)
                results.append({"status": "error", "window": window_title, "message": error_message})
                continue
            win32gui.SetWindowText(hwnd, new_title)
            success_message = f"Window title changed from {window_title} to {new_title}"
            logging.info(success_message)
            results.append({"status": "success", "window": window_title, "new_title": new_title, "message": success_message})

        return results if len(results) > 1 else results[0]
    
    @staticmethod
    @error_handler
    def get_monitor_info():
        logging.info("Retrieving monitor information")
        monitors = []  
        
        def callback(monitor, dc, rect, data):
            info = win32api.GetMonitorInfo(monitor)
            monitors.append(info)
            return True

        try:
            win32api.EnumDisplayMonitors(None, None)
            hMonitors = win32api.EnumDisplayMonitors()
            for hMonitor, _, _ in hMonitors:
                monitors.append(win32api.GetMonitorInfo(hMonitor))
        except Exception as e:
            error_message = f"Error enumerating monitors: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "message": error_message}
        
        logging.info(f"Monitor information: {monitors}")
        return {"status": "success", "monitors": monitors, "message": "Retrieved monitor information"}
    
    @staticmethod
    @error_handler
    def set_window_transparency(window_titles, transparency_levels):
        if isinstance(window_titles, str):
            window_titles = [window_titles]
        if isinstance(transparency_levels, (int, float)):
            transparency_levels = [transparency_levels] * len(window_titles)
        
        results = []
        for window_title, transparency_level in zip(window_titles, transparency_levels):
            logging.info(f"Setting window transparency: {window_title} to level: {transparency_level}")
            hwnd = win32gui.FindWindow(None, window_title)
            if not hwnd:
                error_message = f"Window not found: {window_title}"
                logging.error(error_message)
                results.append({"status": "error", "window": window_title, "message": error_message})
                continue
            try:
                extended_style_settings = win32gui.GetWindowLong(hwnd, win32con.GWL_EXSTYLE)
                win32gui.SetWindowLong(hwnd, win32con.GWL_EXSTYLE, extended_style_settings | win32con.WS_EX_LAYERED)
                win32gui.SetLayeredWindowAttributes(hwnd, win32api.RGB(0, 0, 0), transparency_level, win32con.LWA_ALPHA)
                success_message = f"Window {window_title} transparency set to {transparency_level}"
                logging.info(success_message)
                results.append({"status": "success", "window": window_title, "transparency": transparency_level, "message": success_message})
            except Exception as e:
                error_message = f"Failed to set window transparency for {window_title}: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "window": window_title, "message": error_message})

        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def toggle_always_on_top(window_titles, enable=True):
        if isinstance(window_titles, str):
            window_titles = [window_titles]
        
        results = []
        for window_title in window_titles:
            logging.info(f"Toggling always on top for window: {window_title} to {'enable' if enable else 'disable'}")
            hwnd = win32gui.FindWindow(None, window_title)
            if not hwnd:
                error_message = f"Window not found: {window_title}"
                logging.error(error_message)
                results.append({"status": "error", "window": window_title, "message": error_message})
                continue
            try:
                win32gui.SetWindowPos(hwnd, win32con.HWND_TOPMOST if enable else win32con.HWND_NOTOPMOST, 0, 0, 0, 0, 
                                      win32con.SWP_NOMOVE | win32con.SWP_NOSIZE)
                success_message = f"Window {window_title} always on top set to {'enabled' if enable else 'disabled'}"
                logging.info(success_message)
                results.append({"status": "success", "window": window_title, "always_on_top": enable, "message": success_message})
            except Exception as e:
                error_message = f"Failed to toggle always on top for {window_title}: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "window": window_title, "message": error_message})

        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def restore_window(window_titles):
        if isinstance(window_titles, str):
            window_titles = [window_titles]
        
        results = []
        for window_title in window_titles:
            logging.info(f"Restoring window: {window_title}")
            hwnd = win32gui.FindWindow(None, window_title)
            if not hwnd:
                error_message = f"Window not found: {window_title}"
                logging.error(error_message)
                results.append({"status": "error", "window": window_title, "message": error_message})
                continue
            try:
                win32gui.ShowWindow(hwnd, win32con.SW_RESTORE)
                success_message = f"Window {window_title} restored"
                logging.info(success_message)
                results.append({"status": "success", "window": window_title, "message": success_message})
            except Exception as e:
                error_message = f"Failed to restore window {window_title}: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "window": window_title, "message": error_message})

        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def close_all_windows(exclude_titles=None):
        logging.info("Closing all windows")
        if exclude_titles is None:
            exclude_titles = []
        def callback(hwnd, extra):
            if win32gui.IsWindowVisible(hwnd) and win32gui.GetWindowText(hwnd) not in exclude_titles:
                win32gui.PostMessage(hwnd, win32con.WM_CLOSE, 0, 0)
            return True
        win32gui.EnumWindows(callback, None)
        success_message = "Closed all windows"
        logging.info(success_message)
        return {"status": "success", "message": success_message}

    @staticmethod
    @error_handler
    def snap_window(window_titles, positions):
        if isinstance(window_titles, str):
            window_titles = [window_titles]
        if isinstance(positions, str):
            positions = [positions] * len(window_titles)
        
        results = []
        for window_title, position in zip(window_titles, positions):
            logging.info(f"Snapping window: {window_title} to position: {position}")
            hwnd = win32gui.FindWindow(None, window_title)
            if not hwnd:
                error_message = f"Window not found: {window_title}"
                logging.error(error_message)
                results.append({"status": "error", "window": window_title, "message": error_message})
                continue

            screen_width = win32api.GetSystemMetrics(win32con.SM_CXSCREEN)
            screen_height = win32api.GetSystemMetrics(win32con.SM_CYSCREEN)

            positions_dict = {
                "left": (0, 0, screen_width // 2, screen_height),
                "right": (screen_width // 2, 0, screen_width // 2, screen_height),
                "top": (0, 0, screen_width, screen_height // 2),
                "bottom": (0, screen_height // 2, screen_width, screen_height // 2),
            }

            if position not in positions_dict:
                error_message = f"Invalid position: {position}"
                logging.error(error_message)
                results.append({"status": "error", "window": window_title, "message": error_message})
                continue

            left, top, width, height = positions_dict[position]
            win32gui.MoveWindow(hwnd, left, top, width, height, True)
            success_message = f"Window {window_title} snapped to {position}"
            logging.info(success_message)
            results.append({"status": "success", "window": window_title, "position": position, "message": success_message})

        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def get_window_position_size(window_titles):
        if isinstance(window_titles, str):
            window_titles = [window_titles]
        
        results = []
        for window_title in window_titles:
            logging.info(f"Getting position and size of window: {window_title}")
            hwnd = win32gui.FindWindow(None, window_title)
            if not hwnd:
                error_message = f"Window not found: {window_title}"
                logging.error(error_message)
                results.append({"status": "error", "window": window_title, "message": error_message})
                continue
            rect = win32gui.GetWindowRect(hwnd)
            position_size = {
                "left": rect[0],
                "top": rect[1],
                "right": rect[2],
                "bottom": rect[3],
                "width": rect[2] - rect[0],
                "height": rect[3] - rect[1]
            }
            success_message = f"Position and size of window {window_title} retrieved"
            logging.info(success_message)
            results.append({"status": "success", "window": window_title, "position_size": position_size, "message": success_message})

        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def list_child_windows(window_title):
        logging.info(f"Listing child windows of: {window_title}")
        hwnd = win32gui.FindWindow(None, window_title)
        if not hwnd:
            error_message = f"Window not found: {window_title}"
            logging.error(error_message)
            return {"status": "error", "message": error_message}
        child_windows = []
        def callback(child_hwnd, param):
            child_windows.append(win32gui.GetWindowText(child_hwnd))
            return True
        win32gui.EnumChildWindows(hwnd, callback, None)
        success_message = f"Child windows of {window_title} listed"
        logging.info(success_message)
        return {"status": "success", "child_windows": child_windows, "message": success_message}
    
    @staticmethod
    @error_handler
    def toggle_window_visibility(window_titles):
        if isinstance(window_titles, str):
            window_titles = [window_titles]
        
        results = []
        for window_title in window_titles:
            logging.info(f"Toggling visibility of window: {window_title}")
            hwnd = win32gui.FindWindow(None, window_title)
            if not hwnd:
                error_message = f"Window not found: {window_title}"
                logging.error(error_message)
                results.append({"status": "error", "window": window_title, "message": error_message})
                continue
            try:
                style = win32gui.GetWindowLong(hwnd, win32con.GWL_STYLE)
                if style & win32con.WS_VISIBLE:
                    win32gui.ShowWindow(hwnd, win32con.SW_HIDE)
                    action = "hidden"
                else:
                    win32gui.ShowWindow(hwnd, win32con.SW_SHOW)
                    action = "shown"
                success_message = f"Window {window_title} visibility {action}"
                logging.info(success_message)
                results.append({"status": "success", "window": window_title, "visibility": action, "message": success_message})
            except Exception as e:
                error_message = f"Failed to toggle visibility for {window_title}: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "window": window_title, "message": error_message})

        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def move_window_to_next_monitor(window_titles):
        if isinstance(window_titles, str):
            window_titles = [window_titles]
        
        results = []
        for window_title in window_titles:
            logging.info(f"Moving window: {window_title} to the next monitor")
            hwnd = win32gui.FindWindow(None, window_title)
            if not hwnd:
                error_message = f"Window not found: {window_title}"
                logging.error(error_message)
                results.append({"status": "error", "window": window_title, "message": error_message})
                continue
            
            monitor_info = WindowManagementAgent.get_monitor_info()
            if monitor_info["status"] != "success":
                results.append({"status": "error", "window": window_title, "message": "Failed to get monitor information"})
                continue

            monitors = monitor_info["monitors"]
            if len(monitors) < 2:
                message = "Only one monitor detected. Cannot move to next monitor."
                logging.info(message)
                results.append({"status": "info", "window": window_title, "message": message})
                continue

            current_rect = win32gui.GetWindowRect(hwnd)
            current_monitor = win32api.MonitorFromRect(current_rect)

            next_monitor = None
            for i, monitor in enumerate(monitors):
                if monitor['hMonitor'] == current_monitor:
                    next_monitor = monitors[(i + 1) % len(monitors)]
                    break

            if next_monitor:
                work_area = next_monitor['Work']
                
                new_left = work_area[0] + (work_area[2] - work_area[0] - (current_rect[2] - current_rect[0])) // 2
                new_top = work_area[1] + (work_area[3] - work_area[1] - (current_rect[3] - current_rect[1])) // 2

                win32gui.MoveWindow(hwnd, new_left, new_top, current_rect[2] - current_rect[0], current_rect[3] - current_rect[1], True)
                
                success_message = f"Window {window_title} moved to next monitor"
                logging.info(success_message)
                results.append({"status": "success", "window": window_title, "message": success_message})
            else:
                error_message = f"Failed to find next monitor for {window_title}"
                logging.error(error_message)
                results.append({"status": "error", "window": window_title, "message": error_message})

        return results if len(results) > 1 else results[0]

class ClipboardAgent:
    def __init__(self):
        self.clipboard_history = []

    @staticmethod
    @error_handler
    def copy_to_clipboard(text):
        logging.info(f"Copying text to clipboard: {text}")
        pyperclip.copy(text)
        return {"status": "success", "message": "Text copied to clipboard"}

    @staticmethod
    @error_handler
    def paste_from_clipboard():
        logging.info("Pasting text from clipboard")
        text = pyperclip.paste()
        return {"status": "success", "text": text, "message": "Text pasted from clipboard"}

    @staticmethod
    @error_handler
    def clear_clipboard():
        logging.info("Clearing clipboard")
        pyperclip.copy('')
        return {"status": "success", "message": "Clipboard cleared"}

    @error_handler
    def add_to_history(self, text):
        logging.info(f"Adding text to clipboard history: {text}")
        self.clipboard_history.append(text)
        return {"status": "success", "message": "Text added to clipboard history"}

    @error_handler
    def get_clipboard_history(self):
        logging.info("Retrieving clipboard history")
        return {"status": "success", "history": self.clipboard_history, "message": "Clipboard history retrieved"}

    @staticmethod
    @error_handler
    def monitor_clipboard(callback, interval=1.0):
        logging.info("Starting clipboard monitoring")
        previous_text = pyperclip.paste()
        try:
            while True:
                current_text = pyperclip.paste()
                if current_text != previous_text:
                    logging.info(f"Clipboard content changed: {current_text}")
                    callback(current_text)
                    previous_text = current_text
                time.sleep(interval)
        except KeyboardInterrupt:
            logging.info("Clipboard monitoring stopped")
            return {"status": "success", "message": "Clipboard monitoring stopped"}

    @staticmethod
    @error_handler
    def copy_image_to_clipboard(image_path):
        logging.info(f"Copying image to clipboard from path: {image_path}")
        try:
            image = ImageGrab.open(image_path)
            image.show()
            return {"status": "success", "message": "Image copied to clipboard"}
        except Exception as e:
            logging.error(f"Error copying image to clipboard: {str(e)}")
            return {"status": "error", "message": str(e)}

    @staticmethod
    @error_handler
    def paste_image_from_clipboard():
        logging.info("Pasting image from clipboard")
        try:
            image = ImageGrab.grabclipboard()
            if image:
                image_path = "pasted_image.png"
                image.save(image_path)
                return {"status": "success", "image_path": image_path, "message": "Image pasted from clipboard"}
            else:
                return {"status": "error", "message": "No image found in clipboard"}
        except Exception as e:
            logging.error(f"Error pasting image from clipboard: {str(e)}")
            return {"status": "error", "message": str(e)}
        
    @staticmethod
    @error_handler
    def copy_file_to_clipboard(file_path):
        logging.info(f"Copying file to clipboard: {file_path}")
        
        image = Image.open(file_path)
        output = BytesIO()
        image.convert("RGB").save(output, "BMP")
        data = output.getvalue()[14:]
        output.close()
        
        win32clipboard.OpenClipboard()
        win32clipboard.EmptyClipboard()
        win32clipboard.SetClipboardData(win32clipboard.CF_DIB, data)
        win32clipboard.CloseClipboard()
        
        return {"status": "success", "message": "File copied to clipboard"}

    @staticmethod
    @error_handler
    def get_clipboard_format():
        logging.info("Getting clipboard format")
        
        formats = {
            win32clipboard.CF_TEXT: "Text",
            win32clipboard.CF_BITMAP: "Image",
            win32clipboard.CF_HDROP: "File"
        }
        
        win32clipboard.OpenClipboard()
        format = win32clipboard.GetPriorityClipboardFormat(list(formats.keys()))
        win32clipboard.CloseClipboard()
        
        return {"status": "success", "format": formats.get(format, "Unknown"), "message": "Clipboard format retrieved"}

class RegistryAgent:
    @staticmethod
    def _handle_binary_data(value):
        if isinstance(value, bytes):
            try:
                return value.decode('utf-8')
            except UnicodeDecodeError:
                return value.hex()
        elif isinstance(value, list):
            return [RegistryAgent._handle_binary_data(item) for item in value]
        elif isinstance(value, dict):
            return {k: RegistryAgent._handle_binary_data(v) for k, v in value.items()}
        else:
            return value

    @staticmethod
    @error_handler
    def read_registry_value(key, subkey, value_names):
        if isinstance(value_names, str):
            value_names = [value_names]
        
        results = []
        for value_name in value_names:
            logging.info(f"Reading registry value: {value_name} from {subkey}")
            try:
                key_handle = winreg.OpenKey(key, subkey, 0, winreg.KEY_READ)
                value, _ = winreg.QueryValueEx(key_handle, value_name)
                winreg.CloseKey(key_handle)
                value = RegistryAgent._handle_binary_data(value)
                results.append({"status": "success", "value_name": value_name, "value": value, "message": f"Read registry value: {value_name}"})
            except WindowsError as e:
                error_message = f"Failed to read registry value {value_name}: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "value_name": value_name, "message": error_message})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def write_registry_value(key, subkey, value_data):
        if not isinstance(value_data, list):
            value_data = [value_data]
        
        results = []
        for data in value_data:
            value_name, value, value_type = data
            logging.info(f"Writing registry value: {value_name} to {subkey}")
            try:
                key_handle = winreg.CreateKey(key, subkey)
                winreg.SetValueEx(key_handle, value_name, 0, value_type, value)
                winreg.CloseKey(key_handle)
                results.append({"status": "success", "value_name": value_name, "message": f"Wrote registry value: {value_name}"})
            except WindowsError as e:
                error_message = f"Failed to write registry value {value_name}: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "value_name": value_name, "message": error_message})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def delete_registry_value(key, subkey, value_names):
        if isinstance(value_names, str):
            value_names = [value_names]
        
        results = []
        for value_name in value_names:
            logging.info(f"Deleting registry value: {value_name} from {subkey}")
            try:
                key_handle = winreg.OpenKey(key, subkey, 0, winreg.KEY_ALL_ACCESS)
                winreg.DeleteValue(key_handle, value_name)
                winreg.CloseKey(key_handle)
                results.append({"status": "success", "value_name": value_name, "message": f"Deleted registry value: {value_name}"})
            except WindowsError as e:
                error_message = f"Failed to delete registry value {value_name}: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "value_name": value_name, "message": error_message})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def create_registry_key(key, subkeys):
        if isinstance(subkeys, str):
            subkeys = [subkeys]
        
        results = []
        for subkey in subkeys:
            logging.info(f"Creating registry key: {subkey}")
            try:
                key_handle = winreg.CreateKey(key, subkey)
                winreg.CloseKey(key_handle)
                results.append({"status": "success", "subkey": subkey, "message": f"Created registry key: {subkey}"})
            except WindowsError as e:
                error_message = f"Failed to create registry key {subkey}: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "subkey": subkey, "message": error_message})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def delete_registry_key(key, subkeys):
        if isinstance(subkeys, str):
            subkeys = [subkeys]
        
        results = []
        for subkey in subkeys:
            logging.info(f"Deleting registry key: {subkey}")
            try:
                winreg.DeleteKey(key, subkey)
                results.append({"status": "success", "subkey": subkey, "message": f"Deleted registry key: {subkey}"})
            except WindowsError as e:
                error_message = f"Failed to delete registry key {subkey}: {str(e)}"
                logging.error(error_message)
                results.append({"status": "error", "subkey": subkey, "message": error_message})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def list_subkeys(key, subkey):
        logging.info(f"Listing subkeys of: {subkey}")
        try:
            key_handle = winreg.OpenKey(key, subkey, 0, winreg.KEY_READ)
            subkeys = []
            i = 0
            while True:
                try:
                    subkey_name = winreg.EnumKey(key_handle, i)
                    subkeys.append(subkey_name)
                    i += 1
                except WindowsError:
                    break
            winreg.CloseKey(key_handle)
            return {"status": "success", "subkeys": subkeys, "message": f"Listed subkeys of: {subkey}"}
        except WindowsError as e:
            error_message = f"Failed to list subkeys: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "message": error_message}

    @staticmethod
    @error_handler
    def list_values(key, subkey):
        logging.info(f"Listing values of: {subkey}")
        try:
            key_handle = winreg.OpenKey(key, subkey, 0, winreg.KEY_READ)
            values = []
            i = 0
            while True:
                try:
                    name, data, type = winreg.EnumValue(key_handle, i)
                    data = RegistryAgent._handle_binary_data(data)
                    values.append({"name": name, "data": data, "type": type})
                    i += 1
                except WindowsError:
                    break
            winreg.CloseKey(key_handle)
            return {"status": "success", "values": values, "message": f"Listed values of: {subkey}"}
        except WindowsError as e:
            error_message = f"Failed to list values: {str(e)}"
            logging.error(error_message)
            return {"status": "error", "message": error_message}
        
    @staticmethod
    @error_handler
    def export_registry_key(key, subkeys, file_paths):
        if isinstance(subkeys, str):
            subkeys = [subkeys]
            file_paths = [file_paths]
        
        results = []
        for subkey, file_path in zip(subkeys, file_paths):
            logging.info(f"Exporting registry key: {subkey} to {file_path}")
            try:
                command = f'reg export "{key}\\{subkey}" "{file_path}" /y'
                result = subprocess.run(command, shell=True, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
                results.append({"status": "success", "subkey": subkey, "file_path": file_path, "message": f"Exported registry key: {subkey} to {file_path}"})
            except subprocess.CalledProcessError as e:
                error_message = f"Failed to export registry key {subkey}: {str(e.stderr)}"
                logging.error(error_message)
                results.append({"status": "error", "subkey": subkey, "file_path": file_path, "message": error_message})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def import_registry_key(file_paths):
        if isinstance(file_paths, str):
            file_paths = [file_paths]
        
        results = []
        for file_path in file_paths:
            logging.info(f"Importing registry key from {file_path}")
            try:
                command = f'reg import "{file_path}"'
                result = subprocess.run(command, shell=True, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
                results.append({"status": "success", "file_path": file_path, "message": f"Imported registry key from {file_path}"})
            except subprocess.CalledProcessError as e:
                error_message = f"Failed to import registry key from {file_path}: {str(e.stderr)}"
                logging.error(error_message)
                results.append({"status": "error", "file_path": file_path, "message": error_message})
        
        return results if len(results) > 1 else results[0]


    @staticmethod
    @error_handler
    def backup_registry(file_path):
        logging.info(f"Backing up registry to {file_path}")
        try:
            command = f'reg save HKLM "{file_path}\\SYSTEM"'
            result = subprocess.run(command, shell=True, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
            return {"status": "success", "message": f"Registry backup saved to {file_path}"}
        except subprocess.CalledProcessError as e:
            error_message = f"Failed to backup registry: {str(e.stderr)}"
            logging.error(error_message)
            return {"status": "error", "message": error_message}
        
    @staticmethod
    @error_handler
    def check_registry_key_exists(key, subkeys):
        if isinstance(subkeys, str):
            subkeys = [subkeys]
        
        results = []
        for subkey in subkeys:
            logging.info(f"Checking if registry key exists: {subkey}")
            try:
                winreg.OpenKey(key, subkey, 0, winreg.KEY_READ)
                results.append({"status": "success", "subkey": subkey, "exists": True, "message": f"Registry key exists: {subkey}"})
            except WindowsError:
                results.append({"status": "success", "subkey": subkey, "exists": False, "message": f"Registry key does not exist: {subkey}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def list_registry_tree(key, subkey):
        logging.info(f"Listing registry tree for: {subkey}")
        def recurse_registry(key, subkey, tree=None):
            if tree is None:
                tree = {}
            try:
                with winreg.OpenKey(key, subkey, 0, winreg.KEY_READ) as opened_key:
                    i = 0
                    while True:
                        try:
                            name, data, type = winreg.EnumValue(opened_key, i)
                            tree[name] = RegistryAgent._handle_binary_data(data)
                            i += 1
                        except WindowsError:
                            break
                    i = 0
                    while True:
                        try:
                            subkey_name = winreg.EnumKey(opened_key, i)
                            tree[subkey_name] = {}
                            recurse_registry(key, f"{subkey}\\{subkey_name}", tree[subkey_name])
                            i += 1
                        except WindowsError:
                            break
            except WindowsError:
                pass
            return tree

        result = recurse_registry(key, subkey)
        return {"status": "success", "tree": result, "message": f"Registry tree listed for: {subkey}"}

class AudioOperationsAgent:
    @staticmethod
    def _check_file_exists(file_path):
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

    @staticmethod
    def _load_audio(file_path):
        AudioOperationsAgent._check_file_exists(file_path)
        return librosa.load(file_path)

    @staticmethod
    def _save_audio(file_path, data, samplerate):
        sf.write(file_path, data, samplerate)

    @staticmethod
    @error_handler
    def play_audio(file_paths, durations=None):
        import sounddevice as sd
        if isinstance(file_paths, str):
            file_paths = [file_paths]
        if durations is None:
            durations = [None] * len(file_paths)
        elif isinstance(durations, (int, float)):
            durations = [durations]
        
        results = []
        for file_path, duration in zip(file_paths, durations):
            logging.info(f"Playing audio file: {file_path}")
            try:
                AudioOperationsAgent._check_file_exists(file_path)
                data, fs = AudioOperationsAgent._load_audio(file_path)
                
                if duration is not None:
                    duration_samples = int(duration * fs)
                    data = data[:duration_samples]
                
                sd.play(data, fs)
                
                if duration is not None:
                    sd.sleep(int(duration * 1000))
                else:
                    sd.wait()
                
                played_duration = len(data) / fs
                results.append({
                    "status": "success", 
                    "file_path": file_path,
                    "message": f"Audio played: {file_path} (Duration: {played_duration:.2f} seconds)",
                    "played_duration": played_duration
                })
            except Exception as e:
                results.append({"status": "error", "file_path": file_path, "error_message": f"Error playing audio: {str(e)}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def extract_audio_metadata(input_files):
        if isinstance(input_files, str):
            input_files = [input_files]
        
        results = []
        for input_file in input_files:
            logging.info(f"Extracting metadata from {input_file}")
            try:
                AudioOperationsAgent._check_file_exists(input_file)
                y, sr = AudioOperationsAgent._load_audio(input_file)
                duration = librosa.get_duration(y=y, sr=sr)
                rms = librosa.feature.rms(y=y)[0]
                loudness = librosa.feature.rms(y=y)[0].mean()
                metadata = {
                    "duration": duration,
                    "sample_rate": sr,
                    "channels": 2 if y.ndim > 1 else 1,
                    "loudness": float(loudness),
                    "rms_energy": float(rms.mean()),
                }
                results.append({"status": "success", "file_path": input_file, "metadata": metadata, "message": "Extracted audio metadata"})
            except Exception as e:
                results.append({"status": "error", "file_path": input_file, "error_message": f"Error extracting audio metadata: {str(e)}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def detect_silence(input_files, min_silence_len=1000, silence_thresh=-32):
        if isinstance(input_files, str):
            input_files = [input_files]
        
        results = []
        for input_file in input_files:
            logging.info(f"Detecting silence in audio file: {input_file}")
            try:
                AudioOperationsAgent._check_file_exists(input_file)
                y, sr = AudioOperationsAgent._load_audio(input_file)
                intervals = librosa.effects.split(y, top_db=-silence_thresh)
                silence_ranges = []
                for start, end in intervals:
                    if (end - start) / sr * 1000 >= min_silence_len:
                        silence_ranges.append([start / sr * 1000, end / sr * 1000])
                total_duration = librosa.get_duration(y=y, sr=sr) * 1000
                silence_percentage = sum(end - start for start, end in silence_ranges) / total_duration * 100
                results.append({
                    "status": "success",
                    "file_path": input_file,
                    "silence_ranges": silence_ranges,
                    "silence_percentage": silence_percentage,
                    "message": f"Silence detected in audio file. {silence_percentage:.2f}% of the audio is silence."
                })
            except Exception as e:
                results.append({"status": "error", "file_path": input_file, "error_message": f"Error detecting silence: {str(e)}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def adjust_audio_volume(input_files, output_files, volume_changes):
        if isinstance(input_files, str):
            input_files = [input_files]
            output_files = [output_files]
            volume_changes = [volume_changes]
        
        results = []
        for input_file, output_file, volume_change in zip(input_files, output_files, volume_changes):
            logging.info(f"Adjusting audio volume of {input_file} by {volume_change} dB")
            try:
                AudioOperationsAgent._check_file_exists(input_file)
                y, sr = AudioOperationsAgent._load_audio(input_file)
                adjusted_y = y * (10 ** (volume_change / 20))
                AudioOperationsAgent._save_audio(output_file, adjusted_y, sr)
                results.append({"status": "success", "input_file": input_file, "output_file": output_file, "volume_change": volume_change, "message": f"Audio volume adjusted by {volume_change} dB and saved to {output_file}"})
            except Exception as e:
                results.append({"status": "error", "input_file": input_file, "output_file": output_file, "error_message": f"Error adjusting audio volume: {str(e)}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def apply_audio_effect(input_files, output_files, effects, **kwargs):
        if isinstance(input_files, str):
            input_files = [input_files]
            output_files = [output_files]
            effects = [effects]
        
        results = []
        for input_file, output_file, effect in zip(input_files, output_files, effects):
            logging.info(f"Applying {effect} effect to {input_file}")
            try:
                AudioOperationsAgent._check_file_exists(input_file)
                y, sr = AudioOperationsAgent._load_audio(input_file)
                if effect == "fade_in":
                    duration = kwargs.get('duration', 1000)
                    fade_samples = int(duration * sr / 1000)
                    fade = np.linspace(0, 1, fade_samples)
                    y[:fade_samples] *= fade
                elif effect == "fade_out":
                    duration = kwargs.get('duration', 1000)
                    fade_samples = int(duration * sr / 1000)
                    fade = np.linspace(1, 0, fade_samples)
                    y[-fade_samples:] *= fade
                elif effect == "speed_change":
                    speed_factor = kwargs.get('speed_factor', 1.5)
                    y = librosa.effects.time_stretch(y, rate=speed_factor)
                elif effect == "reverse":
                    y = y[::-1]
                else:
                    results.append({"status": "error", "input_file": input_file, "output_file": output_file, "error_message": "Unsupported effect"})
                    continue
                
                AudioOperationsAgent._save_audio(output_file, y, sr)
                results.append({"status": "success", "input_file": input_file, "output_file": output_file, "effect": effect, "message": f"Applied {effect} effect and saved to {output_file}"})
            except Exception as e:
                results.append({"status": "error", "input_file": input_file, "output_file": output_file, "error_message": f"Error applying audio effect: {str(e)}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def adjust_noise_level(input_files, output_files, noise_adjustments):
        if isinstance(input_files, str):
            input_files = [input_files]
            output_files = [output_files]
            noise_adjustments = [noise_adjustments]
        
        results = []
        for input_file, output_file, noise_adjustment in zip(input_files, output_files, noise_adjustments):
            logging.info(f"Adjusting noise level of {input_file} by {noise_adjustment} dB")
            try:
                AudioOperationsAgent._check_file_exists(input_file)
                y, sr = AudioOperationsAgent._load_audio(input_file)
                
                noise = np.random.randn(len(y))
                noise_amplitude = np.abs(y).mean() * 10**(noise_adjustment/20)
                
                if noise_adjustment > 0:
                    y_noise = y + noise_amplitude * noise
                else:
                    y_noise = y - noise_amplitude * noise
                
                y_noise = librosa.util.normalize(y_noise)
                
                AudioOperationsAgent._save_audio(output_file, y_noise, sr)
                results.append({"status": "success", "input_file": input_file, "output_file": output_file, "noise_adjustment": noise_adjustment, "message": f"Noise level adjusted by {noise_adjustment} dB and saved to {output_file}"})
            except Exception as e:
                results.append({"status": "error", "input_file": input_file, "output_file": output_file, "error_message": f"Error adjusting noise level: {str(e)}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def remove_silence(input_files, output_files, min_silence_len=1000, silence_thresh=-32):
        if isinstance(input_files, str):
            input_files = [input_files]
            output_files = [output_files]
        
        results = []
        for input_file, output_file in zip(input_files, output_files):
            logging.info(f"Removing silence from {input_file}")
            try:
                AudioOperationsAgent._check_file_exists(input_file)
                y, sr = AudioOperationsAgent._load_audio(input_file)
                yt, index = librosa.effects.trim(y, top_db=-silence_thresh)
                AudioOperationsAgent._save_audio(output_file, yt, sr)
                results.append({"status": "success", "input_file": input_file, "output_file": output_file, "message": f"Silence removed and saved to {output_file}"})
            except Exception as e:
                results.append({"status": "error", "input_file": input_file, "output_file": output_file, "error_message": f"Error removing silence: {str(e)}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def normalize_audio(input_files, output_files, target_loudness=-14):
        if isinstance(input_files, str):
            input_files = [input_files]
            output_files = [output_files]
        
        results = []
        for input_file, output_file in zip(input_files, output_files):
            logging.info(f"Normalizing audio file: {input_file}")
            try:
                AudioOperationsAgent._check_file_exists(input_file)
                y, sr = AudioOperationsAgent._load_audio(input_file)
                loudness = librosa.feature.rms(y=y)[0].mean()
                gain = 10 ** ((target_loudness - loudness) / 20)
                y_normalized = y * gain
                
                # Prevent clipping
                max_amp = np.max(np.abs(y_normalized))
                if max_amp > 1.0:
                    y_normalized = y_normalized / max_amp
                
                AudioOperationsAgent._save_audio(output_file, y_normalized, sr)
                results.append({"status": "success", "input_file": input_file, "output_file": output_file, "target_loudness": target_loudness, "message": f"Audio normalized to {target_loudness} LUFS and saved to {output_file}"})
            except Exception as e:
                results.append({"status": "error", "input_file": input_file, "output_file": output_file, "error_message": f"Error normalizing audio: {str(e)}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def merge_audio_files(output_file, *input_files):
        logging.info(f"Merging audio files into {output_file}")
        try:
            combined = None
            sr = None
            for file in input_files:
                AudioOperationsAgent._check_file_exists(file)
                y, sr_temp = AudioOperationsAgent._load_audio(file)
                if sr is None:
                    sr = sr_temp
                elif sr != sr_temp:
                    raise ValueError("All input files must have the same sample rate")
                if combined is None:
                    combined = y
                else:
                    combined = np.concatenate((combined, y))
            AudioOperationsAgent._save_audio(output_file, combined, sr)
            return {"status": "success", "input_files": input_files, "output_file": output_file, "message": f"Audio files merged into {output_file}"}
        except Exception as e:
            return {"status": "error", "input_files": input_files, "output_file": output_file, "error_message": f"Error merging audio files: {str(e)}"}

    @staticmethod
    @error_handler
    def split_audio_file(input_file, splits, output_prefix):
        logging.info(f"Splitting audio file {input_file}")
        try:
            AudioOperationsAgent._check_file_exists(input_file)
            y, sr = AudioOperationsAgent._load_audio(input_file)
            total_duration = len(y) / sr
            split_points = [int(split * total_duration * sr) for split in splits]
            split_points = [0] + split_points + [len(y)]
            
            output_files = []
            for i in range(len(split_points) - 1):
                start = split_points[i]
                end = split_points[i+1]
                output_file = f"{output_prefix}_{i+1}.wav"
                AudioOperationsAgent._save_audio(output_file, y[start:end], sr)
                output_files.append(output_file)
            
            return {"status": "success", "message": f"Audio file split into {len(output_files)} parts", "output_files": output_files}
        except Exception as e:
            return {"status": "error", "error_message": f"Error splitting audio file: {str(e)}"}

    @staticmethod
    @error_handler
    def trim_audio(input_files, output_files, start_times, end_times):
        if isinstance(input_files, str):
            input_files = [input_files]
            output_files = [output_files]
            start_times = [start_times]
            end_times = [end_times]
        
        results = []
        for input_file, output_file, start_time, end_time in zip(input_files, output_files, start_times, end_times):
            logging.info(f"Trimming audio file: {input_file}")
            try:
                AudioOperationsAgent._check_file_exists(input_file)
                y, sr = AudioOperationsAgent._load_audio(input_file)
                start_sample = int(start_time * sr)
                end_sample = int(end_time * sr)
                trimmed_y = y[start_sample:end_sample]
                AudioOperationsAgent._save_audio(output_file, trimmed_y, sr)
                results.append({"status": "success", "input_file": input_file, "output_file": output_file, "start_time": start_time, "end_time": end_time, "message": f"Audio trimmed and saved to {output_file}"})
            except Exception as e:
                results.append({"status": "error", "input_file": input_file, "output_file": output_file, "error_message": f"Error trimming audio: {str(e)}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def add_echo(input_files, output_files, delays=0.5, decays=0.5):
        if isinstance(input_files, str):
            input_files = [input_files]
            output_files = [output_files]
            delays = [delays]
            decays = [decays]
        
        results = []
        for input_file, output_file, delay, decay in zip(input_files, output_files, delays, decays):
            logging.info(f"Adding echo to audio file: {input_file}")
            try:
                AudioOperationsAgent._check_file_exists(input_file)
                y, sr = AudioOperationsAgent._load_audio(input_file)
                delay_samples = int(delay * sr)
                echo_filter = np.zeros(delay_samples + 1)
                echo_filter[0] = 1
                echo_filter[-1] = decay
                y_echo = lfilter([1.0], echo_filter, y)
                y_echo = librosa.util.normalize(y_echo)
                AudioOperationsAgent._save_audio(output_file, y_echo, sr)
                results.append({"status": "success", "input_file": input_file, "output_file": output_file, "delay": delay, "decay": decay, "message": f"Echo added and saved to {output_file}"})
            except Exception as e:
                results.append({"status": "error", "input_file": input_file, "output_file": output_file, "error_message": f"Error adding echo: {str(e)}"})
        
        return results if len(results) > 1 else results[0]
        
class SecurityOperationsAgent:
    @staticmethod
    @error_handler
    def encrypt_file(file_paths, passwords):
        if isinstance(file_paths, str):
            file_paths = [file_paths]
        if isinstance(passwords, str):
            passwords = [passwords] * len(file_paths)
        
        results = []
        for file_path, password in zip(file_paths, passwords):
            logging.info(f"Encrypting file: {file_path}")
            try:
                salt = os.urandom(16)
                kdf = PBKDF2HMAC(
                    algorithm=hashes.SHA256(),
                    length=32,
                    salt=salt,
                    iterations=100000,
                )
                key = base64.urlsafe_b64encode(kdf.derive(password.encode()))
                f = Fernet(key)

                with open(file_path, 'rb') as file:
                    data = file.read()

                encrypted_data = f.encrypt(data)

                with open(file_path + '.encrypted', 'wb') as file:
                    file.write(salt + encrypted_data)

                results.append({"status": "success", "file_path": file_path, "message": f"File encrypted: {file_path}.encrypted"})
            except Exception as e:
                results.append({"status": "error", "file_path": file_path, "message": f"Error encrypting file: {str(e)}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def decrypt_file(file_paths, passwords):
        if isinstance(file_paths, str):
            file_paths = [file_paths]
        if isinstance(passwords, str):
            passwords = [passwords] * len(file_paths)
        
        results = []
        for file_path, password in zip(file_paths, passwords):
            logging.info(f"Decrypting file: {file_path}")
            try:
                with open(file_path, 'rb') as file:
                    salt = file.read(16)
                    encrypted_data = file.read()

                kdf = PBKDF2HMAC(
                    algorithm=hashes.SHA256(),
                    length=32,
                    salt=salt,
                    iterations=100000,
                )
                key = base64.urlsafe_b64encode(kdf.derive(password.encode()))
                f = Fernet(key)

                decrypted_data = f.decrypt(encrypted_data)

                with open(file_path[:-10], 'wb') as file:
                    file.write(decrypted_data)

                results.append({"status": "success", "file_path": file_path, "message": f"File decrypted: {file_path[:-10]}"})
            except Exception as e:
                results.append({"status": "error", "file_path": file_path, "message": f"Error decrypting file: {str(e)}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def generate_encryption_key():
        logging.info("Generating encryption key")
        key = Fernet.generate_key()
        return {"status": "success", "key": key.decode(), "message": "Encryption key generated"}

    @staticmethod
    @error_handler
    def encrypt_text(texts, keys):
        if isinstance(texts, str):
            texts = [texts]
        if isinstance(keys, str):
            keys = [keys] * len(texts)
        
        results = []
        for text, key in zip(texts, keys):
            logging.info("Encrypting text")
            try:
                f = Fernet(key.encode())
                encrypted_text = f.encrypt(text.encode())
                results.append({"status": "success", "encrypted_text": encrypted_text.decode(), "message": "Text encrypted"})
            except Exception as e:
                results.append({"status": "error", "message": f"Error encrypting text: {str(e)}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def decrypt_text(encrypted_texts, keys):
        if isinstance(encrypted_texts, str):
            encrypted_texts = [encrypted_texts]
        if isinstance(keys, str):
            keys = [keys] * len(encrypted_texts)
        
        results = []
        for encrypted_text, key in zip(encrypted_texts, keys):
            logging.info("Decrypting text")
            try:
                f = Fernet(key.encode())
                decrypted_text = f.decrypt(encrypted_text.encode()).decode()
                results.append({"status": "success", "decrypted_text": decrypted_text, "message": "Text decrypted"})
            except Exception as e:
                results.append({"status": "error", "message": f"Error decrypting text: {str(e)}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def check_file_integrity(file_paths):
        if isinstance(file_paths, str):
            file_paths = [file_paths]
        
        results = []
        for file_path in file_paths:
            logging.info(f"Checking file integrity: {file_path}")
            try:
                sha256_hash = hashlib.sha256()
                with open(file_path, "rb") as f:
                    for byte_block in iter(lambda: f.read(4096), b""):
                        sha256_hash.update(byte_block)
                file_hash = sha256_hash.hexdigest()
                results.append({"status": "success", "file_path": file_path, "file_hash": file_hash, "message": "File integrity checked"})
            except Exception as e:
                results.append({"status": "error", "file_path": file_path, "message": f"Error checking file integrity: {str(e)}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def verify_file_integrity(file_paths, expected_hashes):
        if isinstance(file_paths, str):
            file_paths = [file_paths]
            expected_hashes = [expected_hashes]
        
        results = []
        for file_path, expected_hash in zip(file_paths, expected_hashes):
            logging.info(f"Verifying file integrity: {file_path}")
            try:
                actual_hash = SecurityOperationsAgent.check_file_integrity(file_path)
                is_valid = actual_hash['file_hash'] == expected_hash
                results.append({"status": "success", "file_path": file_path, "is_valid": is_valid, "message": "File integrity verified"})
            except Exception as e:
                results.append({"status": "error", "file_path": file_path, "message": f"Error verifying file integrity: {str(e)}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def secure_erase_file(file_paths):
        if isinstance(file_paths, str):
            file_paths = [file_paths]
        
        results = []
        for file_path in file_paths:
            logging.info(f"Securely erasing file: {file_path}")
            try:
                with open(file_path, "ba+", buffering=0) as f:
                    length = f.tell()
                with open(file_path, "br+", buffering=0) as f:
                    f.write(os.urandom(length))
                os.remove(file_path)
                results.append({"status": "success", "file_path": file_path, "message": f"File securely erased: {file_path}"})
            except Exception as e:
                results.append({"status": "error", "file_path": file_path, "message": f"Error securely erasing file: {str(e)}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def check_password_strength(passwords):
        if isinstance(passwords, str):
            passwords = [passwords]
        
        results = []
        for password in passwords:
            logging.info("Checking password strength")
            try:
                length = len(password)
                categories = [
                    any(c.islower() for c in password),
                    any(c.isupper() for c in password),
                    any(c.isdigit() for c in password),
                    any(c in string.punctuation for c in password)
                ]
                strength = sum(categories)
                if length >= 12 and strength == 4:
                    result = "Very Strong"
                elif length >= 10 and strength >= 3:
                    result = "Strong"
                elif length >= 8 and strength >= 2:
                    result = "Medium"
                else:
                    result = "Weak"
                results.append({"status": "success", "password": password, "strength": result, "message": f"Password strength: {result}"})
            except Exception as e:
                results.append({"status": "error", "password": password, "message": f"Error checking password strength: {str(e)}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def generate_totp_secret():
        logging.info("Generating TOTP secret")
        secret = pyotp.random_base32()
        return {"status": "success", "secret": secret, "message": "TOTP secret generated"}

    @staticmethod
    @error_handler
    def get_totp_token(secret):
        logging.info("Generating TOTP token")
        totp = pyotp.TOTP(secret)
        token = totp.now()
        return {"status": "success", "token": token, "message": "TOTP token generated"}

    @staticmethod
    @error_handler
    def verify_totp_token(secret, token):
        logging.info("Verifying TOTP token")
        totp = pyotp.TOTP(secret)
        verified = totp.verify(token)
        return {"status": "success", "verified": verified, "message": "TOTP token verified"}

    @staticmethod
    @error_handler
    def create_digital_signature(data, private_key):
        logging.info("Creating digital signature")
        private_key_obj = serialization.load_pem_private_key(private_key.encode(), password=None)
        signature = private_key_obj.sign(
            data.encode(),
            padding.PSS(
                mgf=padding.MGF1(hashes.SHA256()),
                salt_length=padding.PSS.MAX_LENGTH
            ),
            hashes.SHA256()
        )
        return {"status": "success", "signature": base64.b64encode(signature).decode(), "message": "Digital signature created"}

    @staticmethod
    @error_handler
    def verify_digital_signature(data, signature, public_key):
        logging.info("Verifying digital signature")
        public_key_obj = serialization.load_pem_public_key(public_key.encode())
        try:
            public_key_obj.verify(
                base64.b64decode(signature),
                data.encode(),
                padding.PSS(
                    mgf=padding.MGF1(hashes.SHA256()),
                    salt_length=padding.PSS.MAX_LENGTH
                ),
                hashes.SHA256()
            )
            return {"status": "success", "verified": True, "message": "Digital signature verified"}
        except Exception as e:
            return {"status": "error", "verified": False, "message": str(e)}

    @staticmethod
    @error_handler
    def generate_rsa_key_pair(key_size=2048):
        logging.info(f"Generating RSA key pair with key size: {key_size}")
        private_key = rsa.generate_private_key(public_exponent=65537, key_size=key_size)
        private_pem = private_key.private_bytes(
            encoding=serialization.Encoding.PEM,
            format=serialization.PrivateFormat.PKCS8,
            encryption_algorithm=serialization.NoEncryption()
        )
        public_pem = private_key.public_key().public_bytes(
            encoding=serialization.Encoding.PEM,
            format=serialization.PublicFormat.SubjectPublicKeyInfo
        )
        return {"status": "success", "private_key": private_pem.decode(), "public_key": public_pem.decode(), "message": "RSA key pair generated"}

    @staticmethod
    @error_handler
    def generate_secure_password(count=1, length=16):
        results = []
        for _ in range(count):
            logging.info(f"Generating secure password of length {length}")
            try:
                chars = string.ascii_letters + string.digits + string.punctuation
                password = ''.join(secrets.choice(chars) for i in range(length))
                results.append({"status": "success", "password": password, "message": "Secure password generated"})
            except Exception as e:
                results.append({"status": "error", "message": f"Error generating secure password: {str(e)}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def encrypt_directory(directory_paths, passwords):
        if isinstance(directory_paths, str):
            directory_paths = [directory_paths]
        if isinstance(passwords, str):
            passwords = [passwords] * len(directory_paths)
        
        results = []
        for directory_path, password in zip(directory_paths, passwords):
            logging.info(f"Encrypting directory: {directory_path}")
            try:
                for root, dirs, files in os.walk(directory_path):
                    for file in files:
                        file_path = os.path.join(root, file)
                        SecurityOperationsAgent.encrypt_file(file_path, password)
                results.append({"status": "success", "directory_path": directory_path, "message": f"Directory encrypted: {directory_path}"})
            except Exception as e:
                results.append({"status": "error", "directory_path": directory_path, "message": f"Error encrypting directory: {str(e)}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def hash_password(passwords):
        if isinstance(passwords, str):
            passwords = [passwords]
        
        results = []
        for password in passwords:
            logging.info("Hashing password")
            try:
                salt = bcrypt.gensalt()
                hashed = bcrypt.hashpw(password.encode(), salt)
                results.append({"status": "success", "password": password, "hashed_password": hashed.decode(), "message": "Password hashed"})
            except Exception as e:
                results.append({"status": "error", "password": password, "message": f"Error hashing password: {str(e)}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def verify_password(passwords, hashed_passwords):
        if isinstance(passwords, str):
            passwords = [passwords]
        if isinstance(hashed_passwords, str):
            hashed_passwords = [hashed_passwords] * len(passwords)
        
        results = []
        for password, hashed_password in zip(passwords, hashed_passwords):
            logging.info("Verifying password")
            try:
                is_valid = bcrypt.checkpw(password.encode(), hashed_password.encode())
                results.append({"status": "success", "password": password, "is_valid": is_valid, "message": "Password verified"})
            except Exception as e:
                results.append({"status": "error", "password": password, "message": f"Error verifying password: {str(e)}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def rsa_encrypt(public_key_pems, messages):
        if isinstance(public_key_pems, str):
            public_key_pems = [public_key_pems]
        if isinstance(messages, str):
            messages = [messages] * len(public_key_pems)
        
        results = []
        for public_key_pem, message in zip(public_key_pems, messages):
            logging.info("Encrypting message with RSA public key")
            try:
                public_key = serialization.load_pem_public_key(public_key_pem.encode())
                encrypted = public_key.encrypt(
                    message.encode(),
                    padding.OAEP(
                        mgf=padding.MGF1(algorithm=hashes.SHA256()),
                        algorithm=hashes.SHA256(),
                        label=None
                    )
                )
                results.append({"status": "success", "encrypted_message": base64.b64encode(encrypted).decode(), "message": "Message encrypted with RSA"})
            except Exception as e:
                results.append({"status": "error", "message": f"Error encrypting message with RSA: {str(e)}"})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    @error_handler
    def rsa_decrypt(private_key_pems, encrypted_messages):
        if isinstance(private_key_pems, str):
            private_key_pems = [private_key_pems]
        if isinstance(encrypted_messages, str):
            encrypted_messages = [encrypted_messages] * len(private_key_pems)
        
        results = []
        for private_key_pem, encrypted_message in zip(private_key_pems, encrypted_messages):
            logging.info("Decrypting message with RSA private key")
            try:
                private_key = serialization.load_pem_private_key(private_key_pem.encode(), password=None)
                decrypted = private_key.decrypt(
                    base64.b64decode(encrypted_message),
                    padding.OAEP(
                        mgf=padding.MGF1(algorithm=hashes.SHA256()),
                        algorithm=hashes.SHA256(),
                        label=None
                    )
                )
                results.append({"status": "success", "decrypted_message": decrypted.decode(), "message": "Message decrypted with RSA"})
            except Exception as e:
                results.append({"status": "error", "message": f"Error decrypting message with RSA: {str(e)}"})
        
        return results if len(results) > 1 else results[0]

class FinanceAgent:

    def __init__(self, call_limit: int = 5):
        self.last_call_time = 0
        self.call_limit = call_limit
        self.call_interval = 1 / self.call_limit

    def rate_limit(self) -> None:
        current_time = time.time()
        time_since_last_call = current_time - self.last_call_time
        if time_since_last_call < self.call_interval:
            time.sleep(self.call_interval - time_since_last_call)
        self.last_call_time = time.time()
    
    @staticmethod
    def error_handler(func):
        def wrapper(*args, **kwargs):
            try:
                return func(*args, **kwargs)
            except requests.RequestException as e:
                logging.error(f"API request failed: {str(e)}")
                return {"status": "error", "message": f"API request failed: {str(e)}"}
            except Exception as e:
                logging.error(f"An unexpected error occurred: {str(e)}")
                return {"status": "error", "message": str(e)}
        return wrapper
    
    @staticmethod
    def convert_timestamps(obj: Any) -> Any:
        if isinstance(obj, pd.Timestamp):
            return obj.strftime('%Y-%m-%d %H:%M:%S')
        elif isinstance(obj, dict):
            return {str(k): FinanceAgent.convert_timestamps(v) for k, v in obj.items()}
        elif isinstance(obj, list):
            return [FinanceAgent.convert_timestamps(v) for v in obj]
        elif isinstance(obj, pd.Series):
            return FinanceAgent.convert_timestamps(obj.to_dict())
        elif isinstance(obj, pd.DataFrame):
            return FinanceAgent.convert_timestamps(obj.to_dict(orient='index'))
        else:
            return obj

    @error_handler
    def get_financial_statements(self, tickers: Union[str, List[str]], statement_type: str = 'income') -> Dict[str, Any]:
        self.rate_limit()
        if isinstance(tickers, str):
            tickers = [tickers]
        
        results = {}
        for ticker in tickers:
            logging.info(f"Getting {statement_type} statement for ticker: {ticker}")
            stock = yf.Ticker(ticker)
            if statement_type == 'income':
                statement = stock.income_stmt
            elif statement_type == 'balance':
                statement = stock.balance_sheet
            elif statement_type == 'cash':
                statement = stock.cashflow
            else:
                results[ticker] = {"status": "error", "message": "Invalid statement type. Choose 'income', 'balance', or 'cash'."}
                continue
            
            statement_dict = self.convert_timestamps(statement)
            
            try:
                json.dumps(statement_dict)  
                results[ticker] = {"status": "success", "statement": statement_dict, "type": statement_type, "message": "Financial statement retrieved successfully"}
            except TypeError as e:
                logging.error(f"Error serializing financial statement for {ticker}: {str(e)}")
                results[ticker] = {"status": "error", "message": f"Failed to serialize financial statement: {str(e)}"}

        return results if len(results) > 1 else list(results.values())[0]

    @error_handler
    @hashable_lru_cache(maxsize=100)
    def get_stock_price(self, tickers: Union[str, List[str]]) -> Dict[str, Any]:
        self.rate_limit()
        if isinstance(tickers, str):
            tickers = [tickers]
        
        results = {}
        for ticker in tickers:
            logging.info(f"Getting stock price for ticker: {ticker}")
            stock = yf.Ticker(ticker)
            info = stock.info
            price = info.get('regularMarketPrice') or info.get('currentPrice') or info.get('previousClose')
            if price is None:
                results[ticker] = {"status": "error", "message": "Unable to retrieve stock price"}
            else:
                results[ticker] = {"status": "success", "price": price, "message": "Stock price retrieved successfully"}

        return results if len(results) > 1 else list(results.values())[0]

    @error_handler
    def get_historical_data(self, tickers: Union[str, List[str]], period: str = '1mo', interval: str = '1d') -> Dict[str, Any]:
        self.rate_limit()
        if isinstance(tickers, str):
            tickers = [tickers]
        
        results = {}
        for ticker in tickers:
            logging.info(f"Getting historical data for ticker: {ticker} for period: {period} with interval: {interval}")
            stock = yf.Ticker(ticker)
            hist = stock.history(period=period, interval=interval)
            hist_data = self.convert_timestamps(hist.reset_index().to_dict(orient='records'))
            results[ticker] = {"status": "success", "historical_data": hist_data, "message": "Historical data retrieved successfully"}

        return results if len(results) > 1 else list(results.values())[0]

    @error_handler
    @hashable_lru_cache(maxsize=100)
    def get_stock_info(self, tickers: Union[str, List[str]]) -> Dict[str, Any]:
        self.rate_limit()
        if isinstance(tickers, str):
            tickers = [tickers]
        
        results = {}
        for ticker in tickers:
            logging.info(f"Getting stock information for ticker: {ticker}")
            stock = yf.Ticker(ticker)
            info = self.convert_timestamps(stock.info)
            results[ticker] = {"status": "success", "info": info, "message": "Stock information retrieved successfully"}

        return results if len(results) > 1 else list(results.values())[0]

    @error_handler
    def get_market_summary(self) -> Dict[str, Any]:
        self.rate_limit()
        logging.info("Getting market summary")
        indices = ['^GSPC', '^DJI', '^IXIC']
        summary = {}

        for index in indices:
            ticker = yf.Ticker(index)
            try:
                info = ticker.info
                logging.debug(f"Fetched info for {index}: {info}")
                hist = ticker.history(period="5d") 
                logging.debug(f"Fetched history for {index}: {hist}")
                
                if not hist.empty:
                    if len(hist) >= 2:
                        current_price = hist['Close'].iloc[-1]
                        previous_price = hist['Close'].iloc[-2]
                        change = current_price - previous_price
                        change_percent = (change / previous_price) * 100
                    else:
                        current_price = hist['Close'].iloc[-1]
                        change = change_percent = 'N/A'
                else:
                    logging.warning(f"No historical data found for {index}")
                    current_price = change = change_percent = 'N/A'
                
            except Exception as e:
                logging.error(f"Error fetching data for {index}: {str(e)}")
                info = {}
                current_price = change = change_percent = 'N/A'
            
            summary[index] = {
                'name': info.get('shortName', 'N/A'),
                'price': current_price,
                'change': change,
                'change_percent': change_percent
            }

        return {"status": "success", "market_summary": summary, "message": "Market summary retrieved successfully"}


    @error_handler
    def get_crypto_prices(self, cryptos: Union[str, List[str]]) -> Dict[str, Any]:
        self.rate_limit()
        if isinstance(cryptos, str):
            cryptos = [cryptos]
        
        logging.info(f"Getting cryptocurrency prices for: {cryptos}")
        prices = {}
        for crypto in cryptos:
            ticker = yf.Ticker(f"{crypto}-USD")
            info = ticker.info
            price = info.get('regularMarketPrice', 'N/A')
            if price == 'N/A':
                hist = ticker.history(period="1d")
                if not hist.empty:
                    price = hist['Close'].iloc[-1]
            prices[crypto] = price
        return {"status": "success", "crypto_prices": prices, "message": "Cryptocurrency prices retrieved successfully"}

    @error_handler
    def get_forex_rates(self, base_currency: str, target_currencies: Union[str, List[str]]) -> Dict[str, Any]:
        self.rate_limit()
        if isinstance(target_currencies, str):
            target_currencies = [target_currencies]
        
        logging.info(f"Getting forex rates for base currency: {base_currency} against {target_currencies}")
        rates = {}
        for currency in target_currencies:
            ticker = yf.Ticker(f"{base_currency}{currency}=X")
            info = ticker.info
            rate = info.get('regularMarketPrice', 'N/A')
            if rate == 'N/A':
                hist = ticker.history(period="1d")
                if not hist.empty:
                    rate = hist['Close'].iloc[-1]
            rates[currency] = rate
        return {"status": "success", "forex_rates": rates, "message": "Forex rates retrieved successfully"}

    @error_handler
    def calculate_moving_average(self, tickers: Union[str, List[str]], window: int = 20) -> Dict[str, Any]:
        self.rate_limit()
        if isinstance(tickers, str):
            tickers = [tickers]
        
        results = {}
        for ticker in tickers:
            logging.info(f"Calculating {window}-day moving average for ticker: {ticker}")
            stock = yf.Ticker(ticker)
            hist = stock.history(period="1y")
            ma = hist['Close'].rolling(window=window).mean().iloc[-1]
            results[ticker] = {"status": "success", "moving_average": ma, "window": window, "message": "Moving average calculated successfully"}

        return results if len(results) > 1 else list(results.values())[0]

    @error_handler
    def compare_stocks(self, tickers: List[str], metric: str = 'price') -> Dict[str, Any]:
        self.rate_limit()
        logging.info(f"Comparing stocks: {tickers} based on {metric}")
        comparison = {}
        for ticker in tickers:
            stock = yf.Ticker(ticker)
            info = stock.info
            if metric == 'price':
                comparison[ticker] = info.get('regularMarketPrice', None)
            elif metric in info:
                comparison[ticker] = info[metric]
            else:
                comparison[ticker] = None
        return {"status": "success", "comparison": comparison, "metric": metric, "message": "Stock comparison completed successfully"}

    @error_handler
    def convert_currency(self, amounts: Union[float, List[float]], from_currencies: Union[str, List[str]], to_currencies: Union[str, List[str]]) -> Dict[str, Any]:
        self.rate_limit()
        if isinstance(amounts, (int, float)):
            amounts = [amounts]
        if isinstance(from_currencies, str):
            from_currencies = [from_currencies] * len(amounts)
        if isinstance(to_currencies, str):
            to_currencies = [to_currencies] * len(amounts)
        
        results = []
        for amount, from_currency, to_currency in zip(amounts, from_currencies, to_currencies):
            logging.info(f"Converting {amount} from {from_currency} to {to_currency}")
            ticker = yf.Ticker(f"{from_currency}{to_currency}=X")
            info = ticker.info
            rate = info.get('regularMarketPrice')
            if rate is None or rate == 'N/A':
                hist = ticker.history(period="1d")
                if not hist.empty:
                    rate = hist['Close'].iloc[-1]
                else:
                    results.append({"status": "error", "from": from_currency, "to": to_currency, "amount": amount, "message": "Unable to retrieve exchange rate"})
                    continue
            converted_amount = amount * rate
            results.append({"status": "success", "from": from_currency, "to": to_currency, "amount": amount, "converted_amount": converted_amount, "rate": rate, "message": "Currency conversion completed successfully"})
        
        return results if len(results) > 1 else results[0]
    
    @error_handler
    def get_earnings_calendar(self, tickers: Union[str, List[str]]) -> Dict[str, Any]:
        self.rate_limit()
        if isinstance(tickers, str):
            tickers = [tickers]
        
        results = {}
        for ticker in tickers:
            logging.info(f"Getting earnings calendar for ticker: {ticker}")
            stock = yf.Ticker(ticker)
            calendar = stock.calendar
            if calendar is None:
                results[ticker] = {"status": "error", "message": "No earnings calendar data available"}
            else:
                if isinstance(calendar, dict):
                    calendar_dict = self.convert_timestamps(calendar)
                else:
                    calendar_dict = self.convert_timestamps(calendar.to_dict())
                results[ticker] = {"status": "success", "calendar": calendar_dict, "message": "Earnings calendar retrieved successfully"}

        return results if len(results) > 1 else list(results.values())[0]

    @error_handler
    def get_sector_performance(self) -> Dict[str, Any]:
        self.rate_limit()
        logging.info("Getting sector performance")
        sectors = ["XLY", "XLP", "XLE", "XLF", "XLV", "XLI", "XLB", "XLRE", "XLK", "XLC", "XLU"]
        performance = {}
        for sector in sectors:
            ticker = yf.Ticker(sector)
            info = ticker.info
            performance[sector] = {
                'name': info.get('shortName', ''),
                'change': info.get('regularMarketChangePercent', '')
            }
        return {"status": "success", "sector_performance": performance, "message": "Sector performance retrieved successfully"}
    
class WeatherAgent:
    WEATHER_BASE_URL = "https://api.open-meteo.com/v1/"
    AIR_QUALITY_BASE_URL = "https://air-quality-api.open-meteo.com/v1/"
    GEOCODING_BASE_URL = "https://geocoding-api.open-meteo.com/v1/search"
    
    def __init__(self):
        pass

    @staticmethod
    def error_handler(func):
        def wrapper(*args, **kwargs):
            try:
                return func(*args, **kwargs)
            except requests.RequestException as e:
                logging.error(f"API request failed: {str(e)}")
                return {"status": "error", "message": f"API request failed: {str(e)}"}
            except Exception as e:
                logging.error(f"An unexpected error occurred: {str(e)}")
                return {"status": "error", "message": str(e)}
        return wrapper

    @hashable_lru_cache(maxsize=100)
    @error_handler
    def get_weather(self, latitudes, longitudes):
        if isinstance(latitudes, (int, float)):
            latitudes = [latitudes]
        if isinstance(longitudes, (int, float)):
            longitudes = [longitudes]
        
        results = []
        for lat, lon in zip(latitudes, longitudes):
            logging.info(f"Getting current weather for coordinates: ({lat}, {lon})")
            params = {
                "latitude": lat,
                "longitude": lon,
                "current": "temperature_2m,relative_humidity_2m,wind_speed_10m,weather_code",
                "hourly": "uv_index",
                "timezone": "auto"
            }
            response = requests.get(f"{self.WEATHER_BASE_URL}forecast", params=params)
            response.raise_for_status()
            data = response.json()
            weather_info = {
                "temperature": data["current"]["temperature_2m"],
                "humidity": data["current"]["relative_humidity_2m"],
                "wind_speed": data["current"]["wind_speed_10m"],
                "weather_code": data["current"]["weather_code"],
                "description": self._get_weather_description(data["current"]["weather_code"]),
                "uv_index": data["hourly"]["uv_index"][0] 
            }
            results.append({"status": "success", "coordinates": (lat, lon), "weather": weather_info, "message": "Current weather retrieved successfully"})
        
        return results if len(results) > 1 else results[0]

    @error_handler
    def get_weather_forecast(self, latitudes, longitudes, days=3):
        if isinstance(latitudes, (int, float)):
            latitudes = [latitudes]
        if isinstance(longitudes, (int, float)):
            longitudes = [longitudes]
        
        results = []
        for lat, lon in zip(latitudes, longitudes):
            logging.info(f"Getting weather forecast for coordinates: ({lat}, {lon}) for {days} days")
            params = {
                "latitude": lat,
                "longitude": lon,
                "daily": "temperature_2m_max,temperature_2m_min,precipitation_sum,wind_speed_10m_max,weather_code",
                "hourly": "temperature_2m,precipitation_probability,wind_speed_10m",
                "timezone": "auto",
                "forecast_days": days
            }
            response = requests.get(f"{self.WEATHER_BASE_URL}forecast", params=params)
            response.raise_for_status()
            data = response.json()
            forecast_data = []
            for i in range(days):
                forecast_data.append({
                    "date": data["daily"]["time"][i],
                    "temperature_max": data["daily"]["temperature_2m_max"][i],
                    "temperature_min": data["daily"]["temperature_2m_min"][i],
                    "precipitation": data["daily"]["precipitation_sum"][i],
                    "wind_speed": data["daily"]["wind_speed_10m_max"][i],
                    "weather_code": data["daily"]["weather_code"][i],
                    "description": self._get_weather_description(data["daily"]["weather_code"][i]),
                    "hourly_data": [
                        {
                            "time": data["hourly"]["time"][j],
                            "temperature": data["hourly"]["temperature_2m"][j],
                            "precipitation_probability": data["hourly"]["precipitation_probability"][j],
                            "wind_speed": data["hourly"]["wind_speed_10m"][j]
                        } for j in range(i*24, (i+1)*24)
                    ]
                })
            results.append({"status": "success", "coordinates": (lat, lon), "forecast": forecast_data, "message": "Weather forecast retrieved successfully"})
        
        return results if len(results) > 1 else results[0]

    @error_handler
    def get_historical_weather(self, latitudes, longitudes, start_dates, end_dates):
        if isinstance(latitudes, (int, float)):
            latitudes = [latitudes]
        if isinstance(longitudes, (int, float)):
            longitudes = [longitudes]
        if isinstance(start_dates, str):
            start_dates = [start_dates] * len(latitudes)
        if isinstance(end_dates, str):
            end_dates = [end_dates] * len(latitudes)
        
        results = []
        for lat, lon, start_date, end_date in zip(latitudes, longitudes, start_dates, end_dates):
            logging.info(f"Getting historical weather for coordinates: ({lat}, {lon}) from {start_date} to {end_date}")
            params = {
                "latitude": lat,
                "longitude": lon,
                "start_date": start_date,
                "end_date": end_date,
                "daily": "temperature_2m_max,temperature_2m_min,precipitation_sum,wind_speed_10m_max,weather_code",
                "timezone": "auto"
            }
            response = requests.get(f"{self.WEATHER_BASE_URL}forecast", params=params)
            response.raise_for_status()
            data = response.json()
            historical_data = []
            for i in range(len(data["daily"]["time"])):
                historical_data.append({
                    "date": data["daily"]["time"][i],
                    "temperature_max": data["daily"]["temperature_2m_max"][i],
                    "temperature_min": data["daily"]["temperature_2m_min"][i],
                    "precipitation": data["daily"]["precipitation_sum"][i],
                    "wind_speed": data["daily"]["wind_speed_10m_max"][i],
                    "weather_code": data["daily"]["weather_code"][i],
                    "description": self._get_weather_description(data["daily"]["weather_code"][i])
                })
            results.append({"status": "success", "coordinates": (lat, lon), "historical_data": historical_data, "message": "Historical weather retrieved successfully"})
        
        return results if len(results) > 1 else results[0]

    @error_handler
    def get_weather_multiple_locations(self, locations):
        results = {}
        for location in locations:
            results[f"{location[0]},{location[1]}"] = self.get_weather(location[0], location[1])
        return {"status": "success", "results": results, "message": "Weather retrieved for multiple locations"}

    @staticmethod
    def convert_temperature(temps, from_unit, to_unit):
        if isinstance(temps, (int, float)):
            temps = [temps]
        
        results = []
        for temp in temps:
            if from_unit == to_unit:
                converted_temp = temp
            elif from_unit == "C" and to_unit == "F":
                converted_temp = (temp * 9/5) + 32
            elif from_unit == "F" and to_unit == "C":
                converted_temp = (temp - 32) * 5/9
            else:
                results.append({"status": "error", "message": "Unsupported temperature units. Use 'C' or 'F'."})
                continue
            results.append({"status": "success", "original_temp": temp, "converted_temp": converted_temp, "from_unit": from_unit, "to_unit": to_unit})
        
        return results if len(results) > 1 else results[0]

    @staticmethod
    def _get_weather_description(code):
        weather_codes = {
            0: "Clear sky",
            1: "Mainly clear", 2: "Partly cloudy", 3: "Overcast",
            45: "Fog", 48: "Depositing rime fog",
            51: "Light drizzle", 53: "Moderate drizzle", 55: "Dense drizzle",
            56: "Light freezing drizzle", 57: "Dense freezing drizzle",
            61: "Slight rain", 63: "Moderate rain", 65: "Heavy rain",
            66: "Light freezing rain", 67: "Heavy freezing rain",
            71: "Slight snow fall", 73: "Moderate snow fall", 75: "Heavy snow fall",
            77: "Snow grains",
            80: "Slight rain showers", 81: "Moderate rain showers", 82: "Violent rain showers",
            85: "Slight snow showers", 86: "Heavy snow showers",
            95: "Thunderstorm", 96: "Thunderstorm with slight hail", 99: "Thunderstorm with heavy hail"
        }
        return weather_codes.get(code, "Unknown")
    
    @error_handler
    def get_air_quality(self, latitudes, longitudes):
        if isinstance(latitudes, (int, float)):
            latitudes = [latitudes]
        if isinstance(longitudes, (int, float)):
            longitudes = [longitudes]
        
        results = []
        for lat, lon in zip(latitudes, longitudes):
            logging.info(f"Getting air quality for coordinates: ({lat}, {lon})")
            params = {
                "latitude": lat,
                "longitude": lon,
                "hourly": ["pm10", "pm2_5", "carbon_monoxide", "nitrogen_dioxide", "sulphur_dioxide", "ozone", "aerosol_optical_depth", "dust", "uv_index", "ammonia"],
                "timezone": "auto"
            }
            response = requests.get(f"{self.AIR_QUALITY_BASE_URL}air-quality", params=params)
            response.raise_for_status()
            data = response.json()

            hourly_data = pd.DataFrame({
                "date": pd.date_range(
                    start=pd.to_datetime(data["hourly"]["time"][0]),
                    periods=len(data["hourly"]["time"]),
                    freq="H"
                ),
                **{key: data["hourly"][key] for key in params["hourly"]}
            })

            current_hour = datetime.datetime.now().replace(minute=0, second=0, microsecond=0)
            current_data = hourly_data[hourly_data["date"] == current_hour].iloc[0].to_dict()

            results.append({
                "status": "success",
                "coordinates": (lat, lon),
                "air_quality": {
                    "current": current_data,
                    "hourly": hourly_data.to_dict(orient="records")
                },
                "message": "Air quality retrieved successfully"
            })
        
        return results if len(results) > 1 else results[0]
   
    @error_handler
    def search_city(self, names):
        if isinstance(names, str):
            names = [names]
        
        results = []
        for name in names:
            logging.info(f"Searching for city: {name}")
            
            if not name or not name.strip():
                results.append({"status": "error", "city": name, "message": "City name cannot be empty"})
                continue
            if len(name) > 100:
                results.append({"status": "error", "city": name, "message": "City name is too long"})
                continue
            
            params = {
                "name": name,
                "count": 1,
                "language": "en",
                "format": "json"
            }
            response = requests.get(self.GEOCODING_BASE_URL, params=params)
            response.raise_for_status()
            data = response.json()
            
            if "results" not in data or not data["results"]:
                results.append({"status": "error", "city": name, "message": f"No city found for the name: {name}"})
                continue
            
            city = data["results"][0]
            results.append({
                "status": "success",
                "city": city["name"],
                "country": city["country"],
                "latitude": city["latitude"],
                "longitude": city["longitude"],
                "elevation": city.get("elevation"),
                "timezone": city.get("timezone"),
                "population": city.get("population"),
                "country_code": city.get("country_code"),
                "admin1": city.get("admin1"),
                "admin2": city.get("admin2"),
                "admin3": city.get("admin3"),
                "admin4": city.get("admin4"),
                "message": "City found"
            })
        
        return results if len(results) > 1 else results[0]

        
class VirtualMachineAgent:

    @staticmethod
    @error_handler
    def start_vm(vm_name):
        import virtualbox
        vbox = virtualbox.VirtualBox()
        logging.info(f"Starting virtual machine: {vm_name}")
        try:
            machine = VirtualMachineAgent.vbox.find_machine(vm_name)
            session = virtualbox.Session()
            progress = machine.launch_vm_process(session, "gui", [])
            progress.wait_for_completion()
            return {"status": "success", "message": f"Virtual machine {vm_name} started"}
        except Exception as e:
            logging.error(f"Error starting virtual machine: {str(e)}")
            return {"status": "error", "message": str(e)}

    @staticmethod
    @error_handler
    def stop_vm(vm_name):
        import virtualbox
        vbox = virtualbox.VirtualBox()
        logging.info(f"Stopping virtual machine: {vm_name}")
        try:
            machine = VirtualMachineAgent.vbox.find_machine(vm_name)
            session = machine.create_session()
            session.console.power_down()
            return {"status": "success", "message": f"Virtual machine {vm_name} stopped"}
        except Exception as e:
            logging.error(f"Error stopping virtual machine: {str(e)}")
            return {"status": "error", "message": str(e)}

    @staticmethod
    @error_handler
    def pause_vm(vm_name):
        import virtualbox
        vbox = virtualbox.VirtualBox()
        logging.info(f"Pausing virtual machine: {vm_name}")
        try:
            machine = VirtualMachineAgent.vbox.find_machine(vm_name)
            session = machine.create_session()
            session.console.pause()
            return {"status": "success", "message": f"Virtual machine {vm_name} paused"}
        except Exception as e:
            logging.error(f"Error pausing virtual machine: {str(e)}")
            return {"status": "error", "message": str(e)}

    @staticmethod
    @error_handler
    def resume_vm(vm_name):
        import virtualbox
        vbox = virtualbox.VirtualBox()
        logging.info(f"Resuming virtual machine: {vm_name}")
        try:
            machine = VirtualMachineAgent.vbox.find_machine(vm_name)
            session = machine.create_session()
            session.console.resume()
            return {"status": "success", "message": f"Virtual machine {vm_name} resumed"}
        except Exception as e:
            logging.error(f"Error resuming virtual machine: {str(e)}")
            return {"status": "error", "message": str(e)}

    @staticmethod
    @error_handler
    def list_vms():
        import virtualbox
        vbox = virtualbox.VirtualBox()
        logging.info("Listing all virtual machines")
        try:
            vms = [vm.name for vm in VirtualMachineAgent.vbox.machines]
            return {"status": "success", "vms": vms, "message": "Virtual machines listed successfully"}
        except Exception as e:
            logging.error(f"Error listing virtual machines: {str(e)}")
            return {"status": "error", "message": str(e)}

    @staticmethod
    @error_handler
    def get_vm_info(vm_name):
        import virtualbox
        vbox = virtualbox.VirtualBox()
        logging.info(f"Getting information for virtual machine: {vm_name}")
        try:
            machine = VirtualMachineAgent.vbox.find_machine(vm_name)
            info = {
                "name": machine.name,
                "state": machine.state,
                "os_type": machine.os_type_id,
                "memory_size": machine.memory_size,
                "cpu_count": machine.cpu_count
            }
            return {"status": "success", "vm_info": info, "message": "Virtual machine information retrieved successfully"}
        except Exception as e:
            logging.error(f"Error getting virtual machine information: {str(e)}")
            return {"status": "error", "message": str(e)}

    @staticmethod
    @error_handler
    def take_snapshot(vm_name, snapshot_name):
        import virtualbox
        vbox = virtualbox.VirtualBox()
        logging.info(f"Taking snapshot for virtual machine: {vm_name} with snapshot name: {snapshot_name}")
        try:
            machine = VirtualMachineAgent.vbox.find_machine(vm_name)
            session = machine.create_session()
            progress = session.console.take_snapshot(snapshot_name, "", False)
            progress.wait_for_completion()
            return {"status": "success", "message": f"Snapshot {snapshot_name} taken for virtual machine {vm_name}"}
        except Exception as e:
            logging.error(f"Error taking snapshot: {str(e)}")
            return {"status": "error", "message": str(e)}

    @staticmethod
    @error_handler
    def restore_snapshot(vm_name, snapshot_name):
        import virtualbox
        vbox = virtualbox.VirtualBox()
        logging.info(f"Restoring snapshot for virtual machine: {vm_name} with snapshot name: {snapshot_name}")
        try:
            machine = VirtualMachineAgent.vbox.find_machine(vm_name)
            session = machine.create_session()
            snapshot = machine.find_snapshot(snapshot_name)
            progress = session.console.restore_snapshot(snapshot)
            progress.wait_for_completion()
            return {"status": "success", "message": f"Snapshot {snapshot_name} restored for virtual machine {vm_name}"}
        except Exception as e:
            logging.error(f"Error restoring snapshot: {str(e)}")
            return {"status": "error", "message": str(e)}

    @staticmethod
    @error_handler
    def clone_vm(vm_name, clone_name):
        import virtualbox
        vbox = virtualbox.VirtualBox()
        logging.info(f"Cloning virtual machine: {vm_name} to {clone_name}")
        try:
            machine = VirtualMachineAgent.vbox.find_machine(vm_name)
            clone = VirtualMachineAgent.vbox.create_machine("", clone_name, [], "", "")
            machine.clone_to(clone, [virtualbox.library.CloneMode.machine_state], [])
            VirtualMachineAgent.vbox.register_machine(clone)
            return {"status": "success", "message": f"Virtual machine {vm_name} cloned to {clone_name}"}
        except Exception as e:
            logging.error(f"Error cloning virtual machine: {str(e)}")
            return {"status": "error", "message": str(e)}

class TeacherAgent:
    def __init__(self):
        pass

    def error_handler(func):
        def wrapper(*args, **kwargs):
            try:
                return func(*args, **kwargs)
            except Exception as e:
                logging.error(f"An error occurred: {str(e)}")
                return {"status": "error", "message": str(e)}
        return wrapper

    @staticmethod
    @error_handler
    def web_search_with_content_return(query, num_results=2):
        logging.info(f"Performing web search for: {query}")
        results = []
        for url in search(query, num_results=num_results):
            try:
                response = requests.get(url)
                response.encoding = 'utf-8'
                soup = BeautifulSoup(response.content, 'html.parser')
                text = ' '.join(p.get_text() for p in soup.find_all('p'))
                results.append({'url': url, 'content': text})
            except Exception as e:
                logging.error(f"Error fetching content from {url}: {str(e)}")
        return {"status": "success", "results": results, "message": "Web search completed"}

    @staticmethod
    @error_handler
    def create_powerpoint(topic, slide_details):
        logging.info(f"Creating PowerPoint for topic: {topic}")
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]

        for slide_info in slide_details:
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            body = slide.placeholders[1]
            title.text = slide_info.get('title', '')
            body.text = slide_info.get('text', '')
            for img in slide_info.get('images', []):
                slide.shapes.add_picture(
                    img['path'], Inches(img['left']), Inches(img['top']),
                    Inches(img['width']), Inches(img['height'])
                )

        pptx_file = f"{topic}.pptx"
        prs.save(pptx_file)
        return {"status": "success", "file": pptx_file, "message": f"PowerPoint for {topic} created"}

    @staticmethod
    @error_handler
    def create_word_document(topic, doc_details):
        logging.info(f"Creating Word document for topic: {topic}")
        doc = Document()
        doc.add_heading(topic, 0)

        for doc_info in doc_details:
            doc.add_heading(doc_info.get('title', ''), level=1)
            doc.add_paragraph(doc_info.get('text', ''))
            for img in doc_info.get('images', []):
                doc.add_picture(img['path'], width=Inches(img['width']))

        docx_file = f"{topic}.docx"
        doc.save(docx_file)
        return {"status": "success", "file": docx_file, "message": f"Word document for {topic} created"}

    @staticmethod
    @error_handler
    def generate_quiz(topic, questions, save_path=None):
        logging.info(f"Generating quiz for topic: {topic}")

        if save_path:
            quiz_data = {
                "topic": topic,
                "questions": questions
            }
            
            script_content = f"""
import customtkinter as ctk
import tkinter as tk
from tkinter import messagebox
import json

# Quiz data
quiz_data = {json.dumps(quiz_data, indent=4)}

class QuizApp(ctk.CTk):
    def __init__(self, topic, questions):
        super().__init__()
        self.title(f"Quiz on {{topic}}")
        self.geometry("600x500")
        self.topic = topic
        self.questions = questions
        self.answers = []
        self.create_widgets()

    def create_widgets(self):
        self.main_frame = ctk.CTkScrollableFrame(self)
        self.main_frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=20)

        for i, question in enumerate(self.questions):
            question_frame = ctk.CTkFrame(self.main_frame)
            question_frame.pack(fill=tk.X, padx=10, pady=10)

            ctk.CTkLabel(question_frame, text=f"Q{{i+1}}: {{question['question']}}", wraplength=500).pack(anchor='w')

            answer_var = tk.StringVar(value="")
            self.answers.append(answer_var)

            for option in question['options']:
                ctk.CTkRadioButton(question_frame, text=option, variable=answer_var, value=option).pack(anchor='w', padx=20)

        submit_button = ctk.CTkButton(self, text="Submit", command=self.submit)
        submit_button.pack(pady=10)

    def submit(self):
        try:
            score = 0
            results = []
            for q, a in zip(self.questions, self.answers):
                user_answer = a.get()
                correct = user_answer == q['correct_answer']
                score += 1 if correct else 0
                results.append(f"Q: {{q['question']}}\\nYour answer: {{user_answer}}\\nCorrect answer: {{q['correct_answer']}}\\n{{'Correct!' if correct else 'Incorrect.'}}")

            result_window = ctk.CTkToplevel(self)
            result_window.title("Quiz Results")
            result_window.geometry("500x400")

            result_frame = ctk.CTkScrollableFrame(result_window)
            result_frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=20)

            ctk.CTkLabel(result_frame, text=f"You scored {{score}} out of {{len(self.questions)}}!", font=("Arial", 16, "bold")).pack(pady=10)

            for result in results:
                ctk.CTkLabel(result_frame, text=result, wraplength=450, justify='left').pack(pady=5)

            self.withdraw()  # Hide the main quiz window
        except Exception as e:
            messagebox.showerror("Error", f"An error occurred: {{str(e)}}")

if __name__ == "__main__":
    app = QuizApp(quiz_data["topic"], quiz_data["questions"])
    app.mainloop()
"""
            script_path = os.path.splitext(save_path)[0] + ".py"
            with open(script_path, 'w') as f:
                f.write(script_content)
            
            logging.info(f"Quiz saved to {save_path}")

        return {"status": "success", "message": f"Quiz for {topic} created"}

    @staticmethod
    @error_handler
    def generate_flashcards(topic, flashcards, save_path=None):
        logging.info(f"Generating flashcards for topic: {topic}")

        if save_path:
            flashcard_data = {
                "topic": topic,
                "flashcards": flashcards
            }
            
            script_content = f"""
import customtkinter as ctk
import tkinter as tk
from tkinter import messagebox
import json

# Flashcard data
flashcard_data = {json.dumps(flashcard_data, indent=4)}

class FlashcardApp(ctk.CTk):
    def __init__(self, topic, flashcards):
        super().__init__()
        self.title(f"Flashcards on {{topic}}")
        self.geometry("500x400")
        self.topic = topic
        self.flashcards = flashcards
        self.current_index = 0
        self.card_side = 'front'
        self.create_widgets()

    def create_widgets(self):
        self.main_frame = ctk.CTkFrame(self)
        self.main_frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=20)

        self.flashcard_frame = ctk.CTkFrame(self.main_frame, width=400, height=200)
        self.flashcard_frame.pack(pady=20)
        self.flashcard_frame.pack_propagate(False)

        self.card_text = ctk.CTkLabel(self.flashcard_frame, text="", wraplength=380)
        self.card_text.pack(expand=True)

        button_frame = ctk.CTkFrame(self.main_frame)
        button_frame.pack(pady=10)

        ctk.CTkButton(button_frame, text="Previous", command=self.previous_card).pack(side=tk.LEFT, padx=5)
        ctk.CTkButton(button_frame, text="Flip", command=self.flip_card).pack(side=tk.LEFT, padx=5)
        ctk.CTkButton(button_frame, text="Next", command=self.next_card).pack(side=tk.LEFT, padx=5)

        self.update_flashcard()

    def update_flashcard(self):
        current_card = self.flashcards[self.current_index]
        if self.card_side == 'front':
            self.card_text.configure(text=current_card['front'])
        else:
            self.card_text.configure(text=current_card['back'])

    def flip_card(self):
        self.card_side = 'back' if self.card_side == 'front' else 'front'
        self.update_flashcard()

    def next_card(self):
        if self.current_index < len(self.flashcards) - 1:
            self.current_index += 1
            self.card_side = 'front'
            self.update_flashcard()
        else:
            messagebox.showinfo("End of Flashcards", "You have reached the end of the flashcards.")

    def previous_card(self):
        if self.current_index > 0:
            self.current_index -= 1
            self.card_side = 'front'
            self.update_flashcard()

if __name__ == "__main__":
    app = FlashcardApp(flashcard_data["topic"], flashcard_data["flashcards"])
    app.mainloop()
"""
            script_path = os.path.splitext(save_path)[0] + ".py"
            with open(script_path, 'w') as f:
                f.write(script_content)
            
            logging.info(f"Flashcards saved to {save_path}")

        return {"status": "success", "message": f"Flashcards for {topic} created"}